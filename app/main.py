"""SHU AI Assistant Streamlit application.

This file contains the complete university assistant backend and Streamlit UI.
It is intentionally documented in a beginner-friendly style for final year
project review: each major section explains what the code is doing and why it
matters for chatbot, RAG, automation, notifications, and deployment workflows.
"""

# =========================
# IMPORTS AND PROJECT DEPENDENCIES SECTION
# =========================
# Standard library imports are built into Python and support files, dates,
# security hashes, background threads, logging, email, URLs, and caching.
import os  # Reads environment variables such as API keys and SMTP settings.
import time  # Adds small delays for crawling and scheduler loops.
import json  # Saves and loads simple local JSON data files.
import hashlib  # Creates stable file fingerprints so duplicate PDFs are skipped.
import tempfile  # Creates temporary files for uploaded documents before processing.
import threading  # Runs the automation scheduler in the background.
import schedule  # Provides a lightweight scheduler helper for recurring jobs.
import smtplib  # Sends email notifications through SMTP.
import logging  # Records system activity and errors in shu_chatbot.log.
import requests  # Downloads pages from the SHU website during crawling.
import boto3  # Connects to AWS services such as Bedrock when configured.
import hmac  # Compares admin credentials safely to avoid timing leaks.

# Streamlit and data visualization imports build the web interface and charts.
import streamlit as st  # Main framework used to render the SHU Assistant UI.
import pandas as pd  # Converts analytics and task data into table/chart format.
import plotly.express as px  # Quickly creates dashboard charts.
import plotly.graph_objects as go  # Builds custom Plotly visualizations when needed.

# Date, typing, and data model helpers keep the code organized and explicit.
from datetime import datetime, timedelta  # Handles timestamps, reminders, and reports.
from email.mime.text import MIMEText  # Builds plain-text or HTML email bodies.
from email.mime.multipart import MIMEMultipart  # Combines email headers and body content.
from typing import List, Dict, Optional, Tuple  # Documents expected variable shapes.
from dataclasses import dataclass, field, asdict  # Creates clean data containers.
from enum import Enum  # Defines fixed choices like departments and notification types.
from urllib.parse import urljoin, urlparse  # Safely combines and checks website URLs.
from pathlib import Path  # Works with file and folder paths across operating systems.
from concurrent.futures import ThreadPoolExecutor, as_completed  # Supports parallel page fetching.
from functools import lru_cache  # Caches repeated results for faster reruns.

# Third-party AI, document, crawling, and scheduling libraries.
from bs4 import BeautifulSoup  # Extracts readable text and links from HTML pages.
from dotenv import load_dotenv  # Loads local .env values during development.
from langchain_community.document_loaders import PyPDFLoader  # Reads PDF pages into LangChain documents.
from langchain_text_splitters import RecursiveCharacterTextSplitter  # Splits long documents into searchable chunks.
from langchain_openai import OpenAIEmbeddings, ChatOpenAI  # Calls OpenAI for embeddings and chat completions.
from langchain_chroma import Chroma  # Stores and searches vector embeddings locally with ChromaDB.
from langchain_core.prompts import ChatPromptTemplate  # Creates reusable prompts for the AI model.
from langchain_core.runnables import RunnablePassthrough  # Passes user questions through LangChain chains.
from langchain_core.documents import Document  # Standard document object used by retrievers.
from langchain_core.output_parsers import StrOutputParser  # Converts model output into plain strings.
from docx import Document as DocxDocument  # Reads Microsoft Word documents.
from pptx import Presentation  # Reads Microsoft PowerPoint slides.
from croniter import croniter  # Calculates next run times from cron expressions.

# Load .env values before Config reads credentials and settings.
load_dotenv()  # Makes local .env configuration available before the app starts.


def _read_secret(name: str, default: str = "") -> str:
    """Read credentials from environment variables first, then Streamlit secrets."""
    value = os.getenv(name)
    if value:
        return value
    try:
        return str(st.secrets.get(name, default))
    except Exception:
        return default

# =========================
# LOGGING CONFIGURATION
# =========================

# Configure logging once so both Streamlit actions and background jobs are traceable.
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('shu_chatbot.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger('SHU_Chatbot')  # Named logger makes log lines easy to filter.


# =========================
# DATA MODELS & ENUMS
# =========================

class Department(Enum):
    """Lists the SHU departments that the assistant can route questions to."""
    COMPUTER_SCIENCE = "Computer Science"
    SOFTWARE_ENGINEERING = "Software Engineering"
    ARTIFICIAL_INTELLIGENCE = "Artificial Intelligence"
    CYBER_SECURITY = "Cyber Security"
    ELECTRICAL_ENGINEERING = "Electrical Engineering"
    MANAGEMENT_SCIENCES = "Management Sciences"
    PHARMACY = "Pharmacy"
    ADMISSIONS = "Admissions"
    EXAMINATION = "Examination"
    STUDENT_AFFAIRS = "Student Affairs"
    FINANCE = "Finance"
    LIBRARY = "Library"
    GENERAL = "General"


class NotificationType(Enum):
    """Lists the notification categories supported by the assistant."""
    TIMETABLE_REMINDER = "timetable_reminder"
    ASSIGNMENT_DEADLINE = "assignment_deadline"
    EXAM_SCHEDULE = "exam_schedule"
    ANNOUNCEMENT = "announcement"
    DOCUMENT_PROCESSED = "document_processed"
    SYSTEM_UPDATE = "system_update"


class QueryPriority(Enum):
    """Represents how urgent or important a student query appears to be."""
    LOW = 1
    MEDIUM = 2
    HIGH = 3
    URGENT = 4


@dataclass
class Student:
    """Stores the student profile information used for personalization and alerts."""
    student_id: str
    name: str
    email: str
    program: str
    semester: int
    section: str
    background: str  # Pre-Engineering or Pre-Medical
    preferences: Dict = field(default_factory=dict)
    notification_settings: Dict = field(default_factory=lambda: {
        "email": True,
        "timetable_reminders": True,
        "assignment_alerts": True,
        "exam_reminders": True,
        "announcements": True
    })


@dataclass
class Notification:
    """Stores one in-app or email notification and its delivery status."""
    id: str
    type: NotificationType
    title: str
    message: str
    recipient_id: str
    created_at: datetime
    scheduled_for: Optional[datetime] = None
    sent: bool = False
    sent_at: Optional[datetime] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class ChatSession:
    """Stores one conversation so analytics and personalization can use history."""
    session_id: str
    student_id: Optional[str]
    started_at: datetime
    messages: List[Dict] = field(default_factory=list)
    department_routed: Optional[str] = None
    resolved: bool = False
    satisfaction_score: Optional[int] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class AutomationTask:
    """Stores scheduler metadata for one automated background task."""
    task_id: str
    name: str
    cron_schedule: str
    last_run: Optional[datetime] = None
    next_run: Optional[datetime] = None
    status: str = "active"
    run_count: int = 0
    error_count: int = 0
    last_error: Optional[str] = None


@dataclass
class AnalyticsEvent:
    """Stores one measurable event such as a query, notification, or sync run."""
    event_type: str
    timestamp: datetime
    user_id: Optional[str]
    data: Dict = field(default_factory=dict)


# =========================
# CONFIGURATION
# =========================

# =========================
# CONFIGURATION AND ENVIRONMENT VARIABLES SECTION
# =========================
# Config collects all values that may change between local development,
# Streamlit Cloud, Docker, and AWS deployment. Keeping them in one class makes
# it easier for students to see where keys, paths, and model names come from.
class Config:
    """Central place for environment variables, model names, paths, and app constants."""
    # OpenAI
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
    OPENAI_MODEL = "gpt-4o"
    EMBEDDING_MODEL = "text-embedding-3-small"

    # Admin login credentials.
    # Configure these in .env or Streamlit secrets:
    # ADMIN_USERNAME=admin
    # ADMIN_PASSWORD=shu123
    ADMIN_USERNAME = _read_secret("ADMIN_USERNAME")
    ADMIN_PASSWORD = _read_secret("ADMIN_PASSWORD")

    # AWS
    AWS_REGION = os.getenv("AWS_REGION", "us-east-1")
    AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY_ID")
    AWS_SECRET_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
    BEDROCK_MODEL_ID = "anthropic.claude-3-sonnet-20240229-v1:0"

    # Email
    SMTP_HOST = os.getenv("SMTP_HOST", "smtp.gmail.com")
    SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
    SMTP_USER = os.getenv("SMTP_USER", "")
    SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "")
    EMAIL_FROM = os.getenv("EMAIL_FROM", "noreply@shu.edu.pk")

    # ChromaDB
    CHROMA_WEB_DIR = "chroma_db_web"
    CHROMA_DOC_DIR = "chroma_db_doc"
    CHROMA_KB_DIR = "chroma_db_kb"
    CHROMA_HISTORY_DIR = "chroma_db_history"
    INCOMING_PDFS_DIR = "incoming_pdfs"
    PROCESSED_FILES_PATH = "processed_files.json"

    # Scraping
    SHU_BASE_URL = "https://shu.edu.pk"
    MAX_PAGES = 150
    REQUEST_DELAY = 0.4
    SYNC_INTERVAL_HOURS = 24

    # Database
    DB_PATH = "shu_chatbot.db"
    ANALYTICS_PATH = "analytics_data.json"
    SESSIONS_PATH = "chat_sessions.json"
    NOTIFICATIONS_PATH = "notifications.json"
    STUDENTS_PATH = "students.json"
    TASKS_PATH = "automation_tasks.json"


# =========================
# AWS BEDROCK INTEGRATION
# =========================

# =========================
# AWS BEDROCK OPTIONAL AI CLIENT SECTION
# =========================
# This optional client shows how the project can connect to AWS Bedrock.
# The app mainly uses OpenAI, but this class documents the AWS deployment path.
class BedrockClient:
    """Amazon Bedrock integration for AI response generation."""

    def __init__(self):
        """Create the object and prepare any clients, folders, or in-memory lists it needs."""
        self.client = None
        self._initialize()

    def _initialize(self):
        """Create the AWS Bedrock runtime client when credentials are available."""
        try:
            if Config.AWS_ACCESS_KEY and Config.AWS_SECRET_KEY:
                self.client = boto3.client(
                    'bedrock-runtime',
                    region_name=Config.AWS_REGION,
                    aws_access_key_id=Config.AWS_ACCESS_KEY,
                    aws_secret_access_key=Config.AWS_SECRET_KEY
                )
                logger.info("Amazon Bedrock client initialized successfully")
            else:
                logger.warning("AWS credentials not found. Bedrock disabled.")
        except Exception as e:
            logger.error(f"Failed to initialize Bedrock: {e}")
            self.client = None

    def generate_response(self, prompt: str, max_tokens: int = 2048) -> Optional[str]:
        """Generate response using Amazon Bedrock Claude model."""
        if not self.client:
            return None

        try:
            body = json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": max_tokens,
                "messages": [
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.1,
            })

            response = self.client.invoke_model(
                modelId=Config.BEDROCK_MODEL_ID,
                body=body,
                contentType="application/json",
                accept="application/json"
            )

            response_body = json.loads(response['body'].read())
            return response_body['content'][0]['text']

        except Exception as e:
            logger.error(f"Bedrock generation failed: {e}")
            return None

    def classify_query(self, query: str) -> Dict:
        """Use Bedrock to classify query intent and department."""
        prompt = f"""Classify this university student query into:
1. Department (one of: Computer Science, Software Engineering, AI, Cyber Security,
   Electrical Engineering, Management Sciences, Pharmacy, Admissions, Examination,
   Student Affairs, Finance, Library, General)
2. Priority (LOW, MEDIUM, HIGH, URGENT)
3. Intent (timetable, grades, admission, course_info, faculty, room_location,
   lab_info, exam, assignment, general, complaint)

Query: "{query}"

Respond in JSON format:
{{"department": "...", "priority": "...", "intent": "...", "confidence": 0.0-1.0}}"""

        response = self.generate_response(prompt, max_tokens=200)
        if response:
            try:
                return json.loads(response)
            except json.JSONDecodeError:
                pass

        return {"department": "General", "priority": "MEDIUM",
                "intent": "general", "confidence": 0.5}


# =========================
# NOTIFICATION ENGINE
# =========================

# =========================
# NOTIFICATION SYSTEM SECTION
# =========================
# The notification engine creates, stores, and sends alerts for timetable
# reminders, assignments, exams, announcements, documents, and system updates.
class NotificationEngine:
    """Handles all notification types: email, in-app, scheduled reminders."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.notifications: List[Notification] = []
        self.email_queue: List[Dict] = []
        self._load_notifications()

    def _load_notifications(self):
        """Load notifications from persistent storage."""
        if os.path.exists(Config.NOTIFICATIONS_PATH):
            try:
                with open(Config.NOTIFICATIONS_PATH, 'r') as f:
                    data = json.load(f)
                    self.notifications = [
                        Notification(
                            id=n['id'], type=NotificationType(n['type']),
                            title=n['title'], message=n['message'],
                            recipient_id=n['recipient_id'],
                            created_at=datetime.fromisoformat(n['created_at']),
                            scheduled_for=datetime.fromisoformat(n['scheduled_for']) if n.get('scheduled_for') else None,
                            sent=n.get('sent', False),
                            sent_at=datetime.fromisoformat(n['sent_at']) if n.get('sent_at') else None,
                            metadata=n.get('metadata', {})
                        ) for n in data
                    ]
            except Exception as e:
                logger.error(f"Error loading notifications: {e}")

    def _save_notifications(self):
        """Persist notifications to storage."""
        try:
            data = []
            for n in self.notifications[-1000:]:  # Keep last 1000
                data.append({
                    'id': n.id,
                    'type': n.type.value,
                    'title': n.title,
                    'message': n.message,
                    'recipient_id': n.recipient_id,
                    'created_at': n.created_at.isoformat(),
                    'scheduled_for': n.scheduled_for.isoformat() if n.scheduled_for else None,
                    'sent': n.sent,
                    'sent_at': n.sent_at.isoformat() if n.sent_at else None,
                    'metadata': n.metadata
                })
            with open(Config.NOTIFICATIONS_PATH, 'w') as f:
                json.dump(data, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving notifications: {e}")

    def create_notification(self, type: NotificationType, title: str,
                           message: str, recipient_id: str,
                           scheduled_for: Optional[datetime] = None,
                           metadata: Dict = None) -> Notification:
        """Create a new notification."""
        notification = Notification(
            id=hashlib.md5(f"{recipient_id}{title}{datetime.now().isoformat()}".encode()).hexdigest()[:12],
            type=type,
            title=title,
            message=message,
            recipient_id=recipient_id,
            created_at=datetime.now(),
            scheduled_for=scheduled_for,
            metadata=metadata or {}
        )
        self.notifications.append(notification)
        self._save_notifications()
        logger.info(f"Notification created: {notification.id} - {title}")
        return notification

    def send_email(self, to_email: str, subject: str, body: str,
                   html_body: Optional[str] = None) -> bool:
        """Send email notification."""
        if not Config.SMTP_USER or not Config.SMTP_PASSWORD:
            logger.warning("SMTP not configured. Email not sent.")
            return False

        try:
            msg = MIMEMultipart('alternative')
            msg['From'] = Config.EMAIL_FROM
            msg['To'] = to_email
            msg['Subject'] = f"[SHU Assistant] {subject}"

            msg.attach(MIMEText(body, 'plain'))
            if html_body:
                msg.attach(MIMEText(html_body, 'html'))

            with smtplib.SMTP(Config.SMTP_HOST, Config.SMTP_PORT) as server:
                server.starttls()
                server.login(Config.SMTP_USER, Config.SMTP_PASSWORD)
                server.send_message(msg)

            logger.info(f"Email sent to {to_email}: {subject}")
            return True

        except Exception as e:
            logger.error(f"Email send failed to {to_email}: {e}")
            return False

    def send_timetable_reminder(self, student: Student, classes_today: List[Dict]):
        """Send daily timetable reminder to student."""
        if not student.notification_settings.get('timetable_reminders'):
            return

        class_list = "\n".join([
            f"  • {c['subject']} — {c['time']} — {c['room']} ({c['teacher']})"
            for c in classes_today
        ])

        message = f"""Good morning {student.name}! 📚

Here's your schedule for today:

{class_list}

Have a productive day at SHU! 🎓
"""

        html_body = f"""
<html>
<body style="font-family: Arial, sans-serif; background: #1a1f2e; color: #e2e8f0; padding: 20px;">
    <div style="max-width: 600px; margin: 0 auto; background: #2d3748; border-radius: 12px; padding: 24px;">
        <h2 style="color: #4f8ef7;">📅 Today's Schedule</h2>
        <p>Good morning <strong>{student.name}</strong>!</p>
        <p>Here are your classes for today:</p>
        <table style="width: 100%; border-collapse: collapse; margin: 16px 0;">
            <tr style="background: #4a5568;">
                <th style="padding: 8px; text-align: left; color: #90cdf4;">Subject</th>
                <th style="padding: 8px; text-align: left; color: #90cdf4;">Time</th>
                <th style="padding: 8px; text-align: left; color: #90cdf4;">Room</th>
                <th style="padding: 8px; text-align: left; color: #90cdf4;">Teacher</th>
            </tr>
            {"".join([f'<tr style="border-bottom: 1px solid #4a5568;"><td style="padding: 8px;">{c["subject"]}</td><td style="padding: 8px;">{c["time"]}</td><td style="padding: 8px;">{c["room"]}</td><td style="padding: 8px;">{c["teacher"]}</td></tr>' for c in classes_today])}
        </table>
        <p style="color: #718096; font-size: 0.9rem;">— SHU AI Assistant</p>
    </div>
</body>
</html>
"""

        self.create_notification(
            type=NotificationType.TIMETABLE_REMINDER,
            title="Today's Class Schedule",
            message=message,
            recipient_id=student.student_id,
            metadata={"classes": classes_today}
        )

        if student.notification_settings.get('email'):
            self.send_email(student.email, "Today's Class Schedule", message, html_body)

    def send_assignment_deadline_alert(self, student: Student,
                                       assignment: Dict, days_remaining: int):
        """Send assignment deadline notification."""
        if not student.notification_settings.get('assignment_alerts'):
            return

        urgency = "🔴 URGENT" if days_remaining <= 1 else "🟡 Reminder" if days_remaining <= 3 else "📝 Upcoming"

        message = f"""{urgency} — Assignment Deadline

{student.name}, you have an upcoming deadline:

📝 Assignment: {assignment['title']}
📚 Subject: {assignment['subject']}
👨‍🏫 Teacher: {assignment['teacher']}
📅 Due Date: {assignment['due_date']}
⏰ Days Remaining: {days_remaining}

Don't forget to submit on time! Good luck! 🍀
"""

        self.create_notification(
            type=NotificationType.ASSIGNMENT_DEADLINE,
            title=f"Assignment Due: {assignment['title']}",
            message=message,
            recipient_id=student.student_id,
            metadata={"assignment": assignment, "days_remaining": days_remaining}
        )

        if student.notification_settings.get('email') and days_remaining <= 3:
            self.send_email(student.email,
                          f"Assignment Due in {days_remaining} day(s): {assignment['title']}",
                          message)

    def send_exam_reminder(self, student: Student, exam: Dict, days_until: int):
        """Send exam schedule notification."""
        if not student.notification_settings.get('exam_reminders'):
            return

        message = f"""📋 Exam Reminder

{student.name}, your exam is approaching:

📝 Subject: {exam['subject']}
📅 Date: {exam['date']}
🕐 Time: {exam['time']}
🏫 Venue: {exam['venue']}
⏰ Days Until Exam: {days_until}

Preparation Tips:
• Review lecture notes and slides
• Practice past papers
• Get adequate rest before the exam

Best of luck! 🎓
"""

        self.create_notification(
            type=NotificationType.EXAM_SCHEDULE,
            title=f"Exam in {days_until} days: {exam['subject']}",
            message=message,
            recipient_id=student.student_id,
            metadata={"exam": exam, "days_until": days_until}
        )

        if student.notification_settings.get('email'):
            self.send_email(student.email,
                          f"Exam Reminder: {exam['subject']} in {days_until} days",
                          message)

    def send_announcement(self, title: str, message: str,
                         target_students: List[Student]):
        """Send university announcement to multiple students."""
        for student in target_students:
            if student.notification_settings.get('announcements'):
                self.create_notification(
                    type=NotificationType.ANNOUNCEMENT,
                    title=title,
                    message=message,
                    recipient_id=student.student_id
                )
                if student.notification_settings.get('email'):
                    self.email_queue.append({
                        'to': student.email,
                        'subject': title,
                        'body': message
                    })

        # Process email queue in batch
        self._process_email_queue()

    def _process_email_queue(self):
        """Process queued emails in batches."""
        while self.email_queue:
            email_data = self.email_queue.pop(0)
            self.send_email(email_data['to'], email_data['subject'], email_data['body'])
            time.sleep(0.1)  # Rate limiting

    def get_pending_notifications(self, student_id: str) -> List[Notification]:
        """Get all unsent notifications for a student."""
        return [n for n in self.notifications
                if n.recipient_id == student_id and not n.sent]

    def mark_as_sent(self, notification_id: str):
        """Mark notification as sent."""
        for n in self.notifications:
            if n.id == notification_id:
                n.sent = True
                n.sent_at = datetime.now()
                break
        self._save_notifications()


# =========================
# QUERY ROUTING ENGINE
# =========================

# =========================
# CHATBOT QUERY ROUTING SECTION
# =========================
# Query routing decides which SHU department should handle a question before
# the RAG engine answers it. This helps analytics and escalation workflows.
class QueryRouter:
    """Intelligent query routing system using keyword analysis + AI classification."""

    DEPARTMENT_KEYWORDS = {
        Department.COMPUTER_SCIENCE: [
            "cs", "computer science", "programming", "algorithm", "data structure",
            "database", "compiler", "operating system", "software engineering",
            "machine learning", "deep learning", "artificial intelligence",
            "cloud computing", "devops", "networks", "oop", "python", "java",
            "foit", "information technology"
        ],
        Department.SOFTWARE_ENGINEERING: [
            "se", "software engineering", "software design", "requirements",
            "software architecture", "testing", "agile", "scrum", "sdlc"
        ],
        Department.ARTIFICIAL_INTELLIGENCE: [
            "ai", "artificial intelligence", "machine learning", "deep learning",
            "neural network", "nlp", "computer vision", "data science"
        ],
        Department.CYBER_SECURITY: [
            "cyber", "security", "hacking", "encryption", "firewall",
            "penetration testing", "ethical hacking", "information security"
        ],
        Department.ADMISSIONS: [
            "admission", "apply", "enrollment", "fee structure", "scholarship",
            "merit", "entry test", "eligibility", "registration", "seat"
        ],
        Department.EXAMINATION: [
            "exam", "result", "grade", "gpa", "cgpa", "transcript",
            "datesheet", "paper", "marks", "pass", "fail", "repeat"
        ],
        Department.STUDENT_AFFAIRS: [
            "society", "event", "sports", "hostel", "counseling",
            "student affairs", "extracurricular", "complaint"
        ],
        Department.FINANCE: [
            "fee", "tuition", "payment", "challan", "refund",
            "financial aid", "installment", "dues"
        ],
        Department.LIBRARY: [
            "library", "book", "journal", "research paper",
            "borrowing", "opac", "reference"
        ],
    }

    DEPARTMENT_CONTACTS = {
        Department.COMPUTER_SCIENCE: {
            "email": "cs@shu.edu.pk",
            "phone": "+92-21-111-SHU-SHU",
            "office": "FOIT Building, 2nd Floor",
            "head": "Prof. Dr. Sheikh Muhammad Munaf (Dean FOIT)",
            "incharge": "Mr. Syed Zahid Badshah"
        },
        Department.ADMISSIONS: {
            "email": "admissions@shu.edu.pk",
            "phone": "+92-21-111-SHU-SHU",
            "portal": "https://admissions.shu.edu.pk",
            "office": "Ground Floor, Main Building"
        },
        Department.EXAMINATION: {
            "email": "examination@shu.edu.pk",
            "phone": "+92-21-111-SHU-SHU",
            "portal": "https://eportal.shu.edu.pk",
            "office": "Ground Floor, Main Building"
        },
        Department.FINANCE: {
            "email": "finance@shu.edu.pk",
            "phone": "+92-21-111-SHU-SHU",
            "office": "Ground Floor, Admin Block"
        },
        Department.STUDENT_AFFAIRS: {
            "email": "studentaffairs@shu.edu.pk",
            "phone": "+92-21-111-SHU-SHU",
            "office": "Student Center"
        },
    }

    def __init__(self, bedrock_client: Optional[BedrockClient] = None):
        """Initialize this component and prepare its dependencies for later use."""
        self.bedrock = bedrock_client
        self.routing_history: List[Dict] = []

    def route_query(self, query: str) -> Dict:
        """Route query to appropriate department with confidence scoring."""
        # Step 1: Keyword-based classification
        keyword_result = self._keyword_classify(query)

        # Step 2: AI-based classification (if Bedrock available)
        ai_result = None
        if self.bedrock:
            ai_result = self.bedrock.classify_query(query)

        # Step 3: Combine results
        final_result = self._combine_classifications(keyword_result, ai_result)

        # Add contact information
        dept = Department(final_result['department'])
        if dept in self.DEPARTMENT_CONTACTS:
            final_result['contact_info'] = self.DEPARTMENT_CONTACTS[dept]

        # Log routing
        self.routing_history.append({
            'query': query,
            'result': final_result,
            'timestamp': datetime.now().isoformat()
        })

        return final_result

    def _keyword_classify(self, query: str) -> Dict:
        """Classify query based on keywords."""
        query_lower = query.lower()
        scores = {}

        for dept, keywords in self.DEPARTMENT_KEYWORDS.items():
            score = sum(1 for kw in keywords if kw in query_lower)
            if score > 0:
                scores[dept] = score

        if not scores:
            return {"department": Department.GENERAL.value,
                    "confidence": 0.3, "method": "keyword_default"}

        best_dept = max(scores, key=scores.get)
        max_score = scores[best_dept]
        confidence = min(max_score / 3.0, 1.0)  # Normalize

        return {"department": best_dept.value,
                "confidence": confidence, "method": "keyword"}

    def _combine_classifications(self, keyword_result: Dict,
                                  ai_result: Optional[Dict]) -> Dict:
        """Combine keyword and AI classifications."""
        if not ai_result:
            return keyword_result

        # If both agree, high confidence
        if keyword_result['department'] == ai_result.get('department'):
            return {
                'department': keyword_result['department'],
                'confidence': max(keyword_result['confidence'],
                                ai_result.get('confidence', 0.5)),
                'priority': ai_result.get('priority', 'MEDIUM'),
                'intent': ai_result.get('intent', 'general'),
                'method': 'combined_agree'
            }

        # If AI has higher confidence, prefer AI
        if ai_result.get('confidence', 0) > keyword_result['confidence']:
            return {**ai_result, 'method': 'ai_preferred'}

        return {**keyword_result,
                'priority': ai_result.get('priority', 'MEDIUM'),
                'intent': ai_result.get('intent', 'general'),
                'method': 'keyword_preferred'}

    def get_escalation_info(self, department: str) -> str:
        """Get escalation/contact info for a department."""
        try:
            dept = Department(department)
            if dept in self.DEPARTMENT_CONTACTS:
                info = self.DEPARTMENT_CONTACTS[dept]
                lines = [f"**{dept.value} Department Contact:**"]
                for key, val in info.items():
                    lines.append(f"  • {key.title()}: {val}")
                return "\n".join(lines)
        except ValueError:
            pass
        return "Contact SHU at https://shu.edu.pk/examination-contact/"


# =========================
# TIMETABLE AUTOMATION ENGINE
# =========================

# =========================
# TIMETABLE AND PDF AUTOMATION SECTION
# =========================
# Timetable automation represents university schedule logic and reminder jobs.
# In a full deployment this can read structured timetable data extracted from PDFs.
class TimetableAutomation:
    """Manages timetable-related automation: reminders, schedule lookup, conflicts."""

    # Complete timetable data structure
    TIMETABLE_DATA = {
        "CS": {
            1: {
                "Pre-Engineering": {
                    "A": [
                        {"subject": "Applied Physics", "room": "Lec Hall 202",
                         "teacher": "Tooba Khan", "day": "Monday", "time": "8:30-10:00"},
                        {"subject": "Functional English", "room": "Lec Hall 202",
                         "teacher": "Mazhar", "day": "Monday", "time": "10:00-11:30"},
                        {"subject": "Calculus and Analytical Geometry", "room": "Lec Hall 201",
                         "teacher": "Uzma Javaid", "day": "Tuesday", "time": "8:30-10:00"},
                        {"subject": "Application of Info. & Comm. Tech.", "room": "Lec Hall 202",
                         "teacher": "Ahsan ul Haq", "day": "Tuesday", "time": "10:00-11:30"},
                        {"subject": "Programming Fundamentals", "room": "Lec Hall 201",
                         "teacher": "Dr. Sheeraz Arif", "day": "Wednesday", "time": "8:30-10:00"},
                        {"subject": "Programming Fundamentals (Lab)", "room": "Lec Hall 309 (Lab)",
                         "teacher": "Mariyam Khan", "day": "Thursday", "time": "8:30-11:30"},
                        {"subject": "Application of ICT (Lab)", "room": "Lec Hall 211 (Lab)",
                         "teacher": "Junaid", "day": "Friday", "time": "8:30-11:30"},
                    ]
                }
            },
            # ... (more semester data would be populated from KB)
        }
    }

    def __init__(self, notification_engine: NotificationEngine):
        """Receive the notification engine so timetable events can alert students."""
        self.notification_engine = notification_engine

    def get_today_schedule(self, student: Student) -> List[Dict]:
        """Get today's classes for a student."""
        today = datetime.now().strftime("%A")
        program_data = self.TIMETABLE_DATA.get(student.program, {})
        semester_data = program_data.get(student.semester, {})
        bg_data = semester_data.get(student.background, {})
        section_data = bg_data.get(student.section, [])

        return [c for c in section_data if c.get('day') == today]

    def send_daily_reminders(self, students: List[Student]):
        """Send daily timetable reminders to all students."""
        for student in students:
            classes = self.get_today_schedule(student)
            if classes:
                self.notification_engine.send_timetable_reminder(student, classes)
                logger.info(f"Timetable reminder sent to {student.student_id}")

    def check_room_availability(self, room: str, day: str, time_slot: str) -> bool:
        """Check if a room is available at a given time."""
        for program_data in self.TIMETABLE_DATA.values():
            for semester_data in program_data.values():
                for bg_data in semester_data.values():
                    for section_data in bg_data.values():
                        for class_info in section_data:
                            if (class_info['room'] == room and
                                class_info['day'] == day and
                                class_info['time'] == time_slot):
                                return False
        return True

    def get_teacher_schedule(self, teacher_name: str) -> List[Dict]:
        """Get full schedule for a teacher."""
        schedule = []
        for program, program_data in self.TIMETABLE_DATA.items():
            for semester, semester_data in program_data.items():
                for bg, bg_data in semester_data.items():
                    for section, section_data in bg_data.items():
                        for class_info in section_data:
                            if teacher_name.lower() in class_info['teacher'].lower():
                                schedule.append({
                                    **class_info,
                                    'program': program,
                                    'semester': semester,
                                    'section': section
                                })
        return schedule


# =========================
# DOCUMENT PROCESSING ENGINE
# =========================

# =========================
# DOCUMENT PROCESSING AND VECTOR EMBEDDINGS SECTION
# =========================
# The document processor turns uploaded PDFs, Word files, and PowerPoints into
# small text chunks, creates OpenAI embeddings, and saves them in ChromaDB.
class DocumentProcessor:
    """Automated document processing for PDFs, Word, PowerPoint files."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)  # Converts text chunks into vectors for similarity search.
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=200
        )
        self.processing_queue: List[Dict] = []
        self.processed_documents: List[Dict] = []

    def _file_hash(self, file_path: str) -> str:
        """Calculate a stable hash so the same upload is indexed only once."""
        digest = hashlib.sha256()
        with open(file_path, "rb") as f:
            for block in iter(lambda: f.read(1024 * 1024), b""):
                digest.update(block)
        return digest.hexdigest()

    def _load_processed_files(self) -> Dict:
        """Read the processed-file registry used to skip duplicate document uploads."""
        if not os.path.exists(Config.PROCESSED_FILES_PATH):
            return {}
        try:
            with open(Config.PROCESSED_FILES_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.warning(f"Could not read processed files index: {e}")
            return {}

    def _save_processed_files(self, processed_files: Dict):
        """Write the processed-file registry back to disk after indexing."""
        with open(Config.PROCESSED_FILES_PATH, "w", encoding="utf-8") as f:
            json.dump(processed_files, f, indent=2)

    def _document_store_exists(self) -> bool:
        """Return True only when an existing Chroma document collection has data."""
        if not os.path.exists(Config.CHROMA_DOC_DIR):
            return False
        try:
            vs = Chroma(
                persist_directory=Config.CHROMA_DOC_DIR,
                embedding_function=self.embeddings
            )
            return vs._collection.count() > 0
        except Exception:
            return False

    def _chunk_metadata(self, chunks: List[Document], file_name: str,
                        file_type: str, uploaded_at: str, source: str,
                        file_hash: str) -> List[Document]:
        """Attach searchable provenance to every stored chunk."""
        for chunk in chunks:
            chunk.metadata.update({
                "file_name": file_name,
                "file_type": file_type,
                "uploaded_at": uploaded_at,
                "source": source,
                "file_hash": file_hash
            })
        return chunks

    def _load_document_vector_store(self) -> Optional[Chroma]:
        """Open the ChromaDB document index that stores uploaded file embeddings."""
        try:
            if self._document_store_exists():
                return Chroma(
                    persist_directory=Config.CHROMA_DOC_DIR,
                    embedding_function=self.embeddings
                )
        except Exception as e:
            logger.error(f"Document vector store load failed: {e}")
        return None

    def get_retriever(self):
        """Load the existing document retriever without changing the stored data."""
        vs = self._load_document_vector_store()
        return vs.as_retriever(search_kwargs={"k": 6}) if vs else None

    def process_document(self, file_path: str, file_name: str,
                        student_id: Optional[str] = None) -> Dict:
        """Process a document and return metadata + retriever."""
        start_time = time.time()
        ext = file_name.rsplit(".", 1)[-1].lower()
        file_hash = self._file_hash(file_path)
        uploaded_at = datetime.now().isoformat()
        source = (
            str(Path(Config.INCOMING_PDFS_DIR) / file_name)
            if Path(file_path).parent.name == Config.INCOMING_PDFS_DIR
            else file_name
        )
        processed_files = self._load_processed_files()

        if file_hash in processed_files:
            vs = self._load_document_vector_store()
            processing_time = time.time() - start_time
            logger.info(f"Skipped duplicate document: {file_name}")
            return {
                "file_name": file_name,
                "file_type": ext,
                "num_pages": 0,
                "num_chunks": 0,
                "processing_time": round(processing_time, 2),
                "summary": "Document was already processed.",
                "retriever": vs.as_retriever(search_kwargs={"k": 6}) if vs else None,
                "vector_store": vs,
                "processed_at": processed_files[file_hash].get("processed_at", uploaded_at),
                "student_id": student_id,
                "skipped": True,
                "reason": "duplicate_file_hash"
            }

        try:
            # Load document based on type
            if ext == "pdf":
                documents = PyPDFLoader(file_path).load()
            elif ext in ("docx", "doc"):
                documents = self._load_docx(file_path)
            elif ext in ("pptx", "ppt"):
                documents = self._load_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")

            if not documents:
                raise ValueError("No readable text found in document.")

            # Split into chunks
            chunks = self.text_splitter.split_documents(documents)
            chunks = self._chunk_metadata(
                chunks=chunks,
                file_name=file_name,
                file_type=ext,
                uploaded_at=uploaded_at,
                source=source,
                file_hash=file_hash
            )

            # Append to the existing Chroma document knowledge base when present.
            os.makedirs(Config.CHROMA_DOC_DIR, exist_ok=True)
            if self._document_store_exists():
                vs = Chroma(
                    persist_directory=Config.CHROMA_DOC_DIR,
                    embedding_function=self.embeddings
                )
                vs.add_documents(chunks)
            else:
                vs = Chroma.from_documents(
                    chunks,
                    embedding=self.embeddings,
                    persist_directory=Config.CHROMA_DOC_DIR
                )
            retriever = vs.as_retriever(search_kwargs={"k": 6})

            # Generate summary
            summary = self._generate_summary(documents[:5])

            processing_time = time.time() - start_time

            result = {
                "file_name": file_name,
                "file_type": ext,
                "num_pages": len(documents),
                "num_chunks": len(chunks),
                "processing_time": round(processing_time, 2),
                "summary": summary,
                "retriever": retriever,
                "vector_store": vs,
                "processed_at": datetime.now().isoformat(),
                "student_id": student_id
            }

            self.processed_documents.append({
                k: v for k, v in result.items()
                if k not in ('retriever', 'vector_store')
            })

            processed_files[file_hash] = {
                "file_name": file_name,
                "file_type": ext,
                "source": source,
                "processed_at": result["processed_at"],
                "uploaded_at": uploaded_at,
                "num_chunks": len(chunks),
                "student_id": student_id
            }
            self._save_processed_files(processed_files)

            logger.info(f"Document processed: {file_name} ({len(chunks)} chunks in {processing_time:.1f}s)")
            return result

        except Exception as e:
            logger.error(f"Document processing failed for {file_name}: {e}")
            raise

    def sync_incoming_pdfs(self) -> Dict:
        """Process any new PDFs dropped into incoming_pdfs."""
        os.makedirs(Config.INCOMING_PDFS_DIR, exist_ok=True)
        pdf_paths = sorted(Path(Config.INCOMING_PDFS_DIR).glob("*.pdf"))
        processed, skipped, failed = [], [], []

        # AWS deployment note:
        # In production, an S3 ObjectCreated event can trigger Lambda or an ECS
        # worker with the S3 bucket/key. That worker downloads the file, calls
        # process_document(), and updates this same ChromaDB collection so new
        # PDF chunks are appended without rebuilding the existing knowledge base.
        for pdf_path in pdf_paths:
            try:
                result = self.process_document(str(pdf_path), pdf_path.name)
                if result.get("skipped"):
                    skipped.append(pdf_path.name)
                else:
                    processed.append({
                        "file_name": pdf_path.name,
                        "chunks": result["num_chunks"]
                    })
            except Exception as e:
                failed.append({"file_name": pdf_path.name, "error": str(e)})
                logger.error(f"Incoming PDF sync failed for {pdf_path.name}: {e}")

        return {
            "processed": processed,
            "skipped": skipped,
            "failed": failed,
            "incoming_dir": Config.INCOMING_PDFS_DIR
        }

    def _load_docx(self, path: str) -> List[Document]:
        """Read paragraphs from a Word file and wrap them as LangChain documents."""
        doc = DocxDocument(path)
        return [
            Document(page_content=p.text.strip(),
                    metadata={"paragraph": i, "source": path})
            for i, p in enumerate(doc.paragraphs) if p.text.strip()
        ]

    def _load_pptx(self, path: str) -> List[Document]:
        """Read text from PowerPoint slides and wrap it as LangChain documents."""
        prs = Presentation(path)
        docs = []
        for n, slide in enumerate(prs.slides, 1):
            lines = [
                para.text.strip()
                for shape in slide.shapes if shape.has_text_frame
                for para in shape.text_frame.paragraphs if para.text.strip()
            ]
            if lines:
                docs.append(Document(
                    page_content="\n".join(lines),
                    metadata={"slide": n, "source": path}
                ))
        return docs

    def _generate_summary(self, documents: List[Document]) -> str:
        """Generate a brief summary of the document."""
        try:
            llm = ChatOpenAI(model=Config.OPENAI_MODEL, temperature=0.1)
            content = "\n".join([d.page_content[:500] for d in documents[:3]])
            prompt = f"Summarize this document in 2-3 sentences:\n\n{content}"
            response = llm.invoke(prompt)
            return response.content
        except Exception:
            return "Document processed successfully."

    def queue_document(self, file_data: Dict):
        """Add document to processing queue."""
        self.processing_queue.append({
            **file_data,
            'queued_at': datetime.now().isoformat(),
            'status': 'pending'
        })

    def process_queue(self):
        """Process all queued documents."""
        while self.processing_queue:
            item = self.processing_queue.pop(0)
            try:
                self.process_document(
                    item['path'], item['name'], item.get('student_id')
                )
                item['status'] = 'completed'
            except Exception as e:
                item['status'] = 'failed'
                item['error'] = str(e)
                logger.error(f"Queue processing failed: {e}")


# =========================
# WEB SYNC ENGINE
# =========================

# =========================
# SHU WEBSITE CRAWLING AND WEB INDEX SECTION
# =========================
# The web sync engine crawls shu.edu.pk, extracts useful page text, splits it
# into chunks, embeds the chunks, and stores them in a searchable ChromaDB index.
class WebSyncEngine:
    """Automated web scraping and synchronization from shu.edu.pk."""

    SHU_SEEDS = [
        "https://shu.edu.pk", "https://shu.edu.pk/about-us/",
        "https://shu.edu.pk/Faculty-of-Engineering/",
        "https://shu.edu.pk/Faculty-of-Information-Technology/",
        "https://shu.edu.pk/Faculty-of-Management-Science/",
        "https://shu.edu.pk/Faculty-of-Science/",
        "https://shu.edu.pk/Faculty-of-Pharmacy/",
        "https://shu.edu.pk/department/department-of-computer-science/",
        "https://shu.edu.pk/admission/",
        "https://shu.edu.pk/admission/admission-process/",
        "https://shu.edu.pk/admission/scholarships-financial-aid/",
        "https://shu.edu.pk/admission/tution/",
        "https://shu.edu.pk/FoIT/bs-computer-science/",
        "https://shu.edu.pk/FoIT/bs-artificial-intelligence/",
        "https://shu.edu.pk/programs/",
        "https://shu.edu.pk/news/",
    ]

    HEADERS = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Accept": "text/html,application/xhtml+xml;q=0.9,*/*;q=0.8",
    }

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)
        self.last_sync: Optional[datetime] = None
        self.sync_history: List[Dict] = []
        self.pages_indexed: int = 0
        self.chunks_indexed: int = 0

    def full_sync(self, progress_callback=None) -> Dict:
        """Perform full website synchronization."""
        start_time = time.time()
        logger.info("Starting full website synchronization...")

        # Crawl
        docs = self._crawl_website(progress_callback)

        if not docs:
            return {"status": "failed", "error": "No content scraped"}

        # Build vector store
        splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
        chunks = splitter.split_documents(docs)

        os.makedirs(Config.CHROMA_WEB_DIR, exist_ok=True)
        vs = Chroma.from_documents(
            chunks, embedding=self.embeddings,
            persist_directory=Config.CHROMA_WEB_DIR
        )

        retriever = vs.as_retriever(search_kwargs={"k": 12})
        sync_time = time.time() - start_time

        self.last_sync = datetime.now()
        self.pages_indexed = len(docs)
        self.chunks_indexed = len(chunks)

        result = {
            "status": "success",
            "pages_scraped": len(docs),
            "chunks_created": len(chunks),
            "sync_time": round(sync_time, 2),
            "timestamp": self.last_sync.isoformat(),
            "retriever": retriever,
            "vector_store": vs
        }

        self.sync_history.append({
            k: v for k, v in result.items()
            if k not in ('retriever', 'vector_store')
        })

        logger.info(f"Sync complete: {len(docs)} pages, {len(chunks)} chunks in {sync_time:.1f}s")
        return result

    def incremental_sync(self) -> Dict:
        """Perform incremental sync (only new/changed pages)."""
        # Check if full sync needed
        if not self.last_sync or (datetime.now() - self.last_sync).hours >= Config.SYNC_INTERVAL_HOURS:
            return self.full_sync()

        # Otherwise just check key pages for updates
        updated_docs = []
        for url in self.SHU_SEEDS[:10]:  # Check seed pages
            soup = self._fetch_page(url)
            if soup:
                content = self._extract_content(soup, url)
                if len(content) >= 60:
                    updated_docs.append(Document(
                        page_content=content,
                        metadata={"source": url, "synced_at": datetime.now().isoformat()}
                    ))

        if updated_docs:
            # Update existing vector store
            try:
                vs = Chroma(persist_directory=Config.CHROMA_WEB_DIR,
                           embedding_function=self.embeddings)
                splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                chunks = splitter.split_documents(updated_docs)
                vs.add_documents(chunks)
                logger.info(f"Incremental sync: {len(chunks)} chunks updated")
                return {"status": "success", "chunks_updated": len(chunks)}
            except Exception as e:
                logger.error(f"Incremental sync failed: {e}")

        return {"status": "no_updates"}

    def _crawl_website(self, progress_callback=None) -> List[Document]:
        """Crawl shu.edu.pk website."""
        visited = set()
        queue = list(self.SHU_SEEDS)
        docs = []
        skip_exts = (".pdf", ".jpg", ".jpeg", ".png", ".gif", ".zip",
                     ".docx", ".xlsx", ".mp4", ".svg", ".ico", ".css", ".js")
        page_count = 0

        while queue and page_count < Config.MAX_PAGES:
            url = queue.pop(0)
            norm = url.split("#")[0].rstrip("/")
            if norm in visited:
                continue
            visited.add(norm)

            soup = self._fetch_page(url)
            if soup is None:
                continue

            content = self._extract_content(soup, url)
            if len(content) >= 60:
                title = soup.find("title")
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": url,
                        "title": title.get_text(strip=True) if title else url,
                        "scraped_at": datetime.now().isoformat()
                    }
                ))

            page_count += 1
            if progress_callback:
                progress_callback(page_count, Config.MAX_PAGES, len(docs))

            # Extract links
            for a in soup.find_all("a", href=True):
                href = a["href"].strip()
                if not href or href.startswith(("mailto:", "tel:", "javascript:", "#")):
                    continue
                full = urljoin(url, href).split("#")[0].rstrip("/")
                parsed = urlparse(full)
                if ("shu.edu.pk" in parsed.netloc
                        and full not in visited and full not in queue
                        and not any(full.lower().endswith(e) for e in skip_exts)):
                    queue.append(full)

            time.sleep(Config.REQUEST_DELAY)

        return docs

    def _fetch_page(self, url: str):
        """Download a single SHU website page for the crawler."""
        try:
            r = requests.get(url, headers=self.HEADERS, timeout=15)
            if r.status_code == 200 and "text/html" in r.headers.get("Content-Type", ""):
                return BeautifulSoup(r.content, "html.parser")
        except Exception:
            pass
        return None

    def _extract_content(self, soup, url: str) -> str:
        """Extract readable text and internal links from one HTML page."""
        import re
        parts = []
        title = soup.find("title")
        if title:
            t = re.sub(r"\s*[-–|]\s*Salim Habib University.*", "",
                      title.get_text(strip=True)).strip()
            if t:
                parts.append(f"PAGE: {t}")

        for tag in soup(["script", "style", "noscript", "iframe", "nav", "header", "footer"]):
            tag.decompose()

        body = soup.get_text("\n", strip=True)
        lines = [l.strip() for l in body.splitlines() if l.strip() and len(l.strip()) > 20]
        seen, unique = set(), []
        for l in lines:
            if l not in seen:
                seen.add(l)
                unique.append(l)
        if unique:
            parts.append("CONTENT:\n" + "\n".join(unique[:140]))

        return "\n\n".join(parts)


# =========================
# CHAT HISTORY & PERSONALIZATION ENGINE
# =========================

# =========================
# CHAT HISTORY AND PERSONALIZATION SECTION
# =========================
# Chat history is saved so the assistant can show previous messages, produce
# suggestions, and support analytics without needing a separate database server.
class ChatHistoryManager:
    """Manages chat sessions, history, and personalization."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.sessions: Dict[str, ChatSession] = {}
        self.embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)
        self._load_sessions()

    def _load_sessions(self):
        """Load sessions from persistent storage."""
        if os.path.exists(Config.SESSIONS_PATH):
            try:
                with open(Config.SESSIONS_PATH, 'r') as f:
                    data = json.load(f)
                    for sid, sdata in data.items():
                        self.sessions[sid] = ChatSession(
                            session_id=sid,
                            student_id=sdata.get('student_id'),
                            started_at=datetime.fromisoformat(sdata['started_at']),
                            messages=sdata.get('messages', []),
                            department_routed=sdata.get('department_routed'),
                            resolved=sdata.get('resolved', False),
                            satisfaction_score=sdata.get('satisfaction_score'),
                            metadata=sdata.get('metadata', {})
                        )
            except Exception as e:
                logger.error(f"Error loading sessions: {e}")

    def _save_sessions(self):
        """Persist sessions to storage."""
        try:
            data = {}
            # Keep only last 100 sessions
            recent_sessions = dict(list(self.sessions.items())[-100:])
            for sid, session in recent_sessions.items():
                data[sid] = {
                    'student_id': session.student_id,
                    'started_at': session.started_at.isoformat(),
                    'messages': session.messages[-50:],  # Keep last 50 messages per session
                    'department_routed': session.department_routed,
                    'resolved': session.resolved,
                    'satisfaction_score': session.satisfaction_score,
                    'metadata': session.metadata
                }
            with open(Config.SESSIONS_PATH, 'w') as f:
                json.dump(data, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving sessions: {e}")

    def create_session(self, student_id: Optional[str] = None) -> str:
        """Create new chat session."""
        session_id = hashlib.md5(
            f"{student_id}{datetime.now().isoformat()}".encode()
        ).hexdigest()[:16]

        self.sessions[session_id] = ChatSession(
            session_id=session_id,
            student_id=student_id,
            started_at=datetime.now()
        )
        self._save_sessions()
        return session_id

    def add_message(self, session_id: str, role: str, content: str,
                   metadata: Dict = None):
        """Add message to session."""
        if session_id not in self.sessions:
            session_id = self.create_session()

        self.sessions[session_id].messages.append({
            'role': role,
            'content': content,
            'timestamp': datetime.now().isoformat(),
            'metadata': metadata or {}
        })
        self._save_sessions()

    def get_session_context(self, session_id: str, last_n: int = 10) -> str:
        """Get recent conversation context for personalization."""
        if session_id not in self.sessions:
            return ""

        messages = self.sessions[session_id].messages[-last_n:]
        context_parts = []
        for msg in messages:
            role = "Student" if msg['role'] == 'user' else "Assistant"
            context_parts.append(f"{role}: {msg['content'][:200]}")

        return "\n".join(context_parts)

    def get_student_history(self, student_id: str) -> List[ChatSession]:
        """Get all sessions for a student."""
        return [s for s in self.sessions.values() if s.student_id == student_id]

    def get_frequent_topics(self, student_id: str) -> List[str]:
        """Analyze student's frequently asked topics."""
        sessions = self.get_student_history(student_id)
        all_messages = []
        for s in sessions:
            all_messages.extend([m['content'] for m in s.messages if m['role'] == 'user'])

        # Simple topic extraction
        topic_keywords = {
            "timetable": ["schedule", "class", "room", "time", "timetable"],
            "grades": ["grade", "gpa", "result", "marks", "score"],
            "teachers": ["teacher", "professor", "dr.", "instructor"],
            "labs": ["lab", "computer lab", "computing"],
            "assignments": ["assignment", "homework", "deadline", "submit"],
            "exams": ["exam", "test", "paper", "midterm", "final"],
        }

        topic_counts = {}
        for msg in all_messages:
            msg_lower = msg.lower()
            for topic, keywords in topic_keywords.items():
                if any(kw in msg_lower for kw in keywords):
                    topic_counts[topic] = topic_counts.get(topic, 0) + 1

        # Sort by frequency
        sorted_topics = sorted(topic_counts.items(), key=lambda x: x[1], reverse=True)
        return [t[0] for t in sorted_topics[:5]]

    def generate_personalized_suggestions(self, student_id: str) -> List[str]:
        """Generate personalized question suggestions based on history."""
        topics = self.get_frequent_topics(student_id)

        suggestions_map = {
            "timetable": "What's my schedule for tomorrow?",
            "grades": "What's the minimum CGPA requirement?",
            "teachers": "Who teaches my next class?",
            "labs": "Which lab is available right now?",
            "assignments": "Any upcoming assignment deadlines?",
            "exams": "When is my next exam?",
        }

        suggestions = [suggestions_map.get(t, f"Tell me about {t}") for t in topics]

        # Add default suggestions if not enough
        defaults = [
            "Where is my next class?",
            "Who is the dean of FOIT?",
            "What are the scholarship options?",
        ]
        while len(suggestions) < 4:
            for d in defaults:
                if d not in suggestions:
                    suggestions.append(d)
                    break
            else:
                break

        return suggestions[:6]


# =========================
# ANALYTICS ENGINE
# =========================

# =========================
# ANALYTICS SECTION
# =========================
# Analytics events feed the admin dashboard with usage, routing, response time,
# satisfaction, and automation information.
class AnalyticsEngine:
    """Tracks and analyzes chatbot usage, performance, and user behavior."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.events: List[AnalyticsEvent] = []
        self._load_events()

    def _load_events(self):
        """Load analytics events from disk so dashboards survive app restarts."""
        if os.path.exists(Config.ANALYTICS_PATH):
            try:
                with open(Config.ANALYTICS_PATH, 'r') as f:
                    data = json.load(f)
                    self.events = [
                        AnalyticsEvent(
                            event_type=e['event_type'],
                            timestamp=datetime.fromisoformat(e['timestamp']),
                            user_id=e.get('user_id'),
                            data=e.get('data', {})
                        ) for e in data[-5000:]  # Keep last 5000
                    ]
            except Exception:
                pass

    def _save_events(self):
        """Persist recent analytics events to the local JSON file."""
        try:
            data = [{
                'event_type': e.event_type,
                'timestamp': e.timestamp.isoformat(),
                'user_id': e.user_id,
                'data': e.data
            } for e in self.events[-5000:]]
            with open(Config.ANALYTICS_PATH, 'w') as f:
                json.dump(data, f)
        except Exception as e:
            logger.error(f"Error saving analytics: {e}")

    def track_event(self, event_type: str, user_id: Optional[str] = None,
                   data: Dict = None):
        """Track an analytics event."""
        event = AnalyticsEvent(
            event_type=event_type,
            timestamp=datetime.now(),
            user_id=user_id,
            data=data or {}
        )
        self.events.append(event)
        self._save_events()

    def get_daily_stats(self, days: int = 7) -> Dict:
        """Get daily usage statistics."""
        cutoff = datetime.now() - timedelta(days=days)
        recent = [e for e in self.events if e.timestamp >= cutoff]

        daily = {}
        for e in recent:
            day = e.timestamp.strftime("%Y-%m-%d")
            if day not in daily:
                daily[day] = {"queries": 0, "sessions": 0, "documents": 0,
                             "unique_users": set()}
            if e.event_type == "query":
                daily[day]["queries"] += 1
            elif e.event_type == "session_start":
                daily[day]["sessions"] += 1
            elif e.event_type == "document_processed":
                daily[day]["documents"] += 1
            if e.user_id:
                daily[day]["unique_users"].add(e.user_id)

        # Convert sets to counts
        for day in daily:
            daily[day]["unique_users"] = len(daily[day]["unique_users"])

        return daily

    def get_department_distribution(self) -> Dict:
        """Get query distribution by department."""
        dept_counts = {}
        for e in self.events:
            if e.event_type == "query_routed":
                dept = e.data.get("department", "General")
                dept_counts[dept] = dept_counts.get(dept, 0) + 1
        return dept_counts

    def get_response_time_stats(self) -> Dict:
        """Get response time statistics."""
        times = [e.data.get("response_time", 0)
                for e in self.events if e.event_type == "response_generated"]
        if not times:
            return {"avg": 0, "min": 0, "max": 0, "p95": 0}

        times.sort()
        return {
            "avg": round(sum(times) / len(times), 2),
            "min": round(min(times), 2),
            "max": round(max(times), 2),
            "p95": round(times[int(len(times) * 0.95)] if times else 0, 2),
            "total_responses": len(times)
        }

    def get_popular_queries(self, top_n: int = 10) -> List[Dict]:
        """Get most frequently asked query types."""
        intent_counts = {}
        for e in self.events:
            if e.event_type == "query":
                intent = e.data.get("intent", "unknown")
                intent_counts[intent] = intent_counts.get(intent, 0) + 1

        sorted_intents = sorted(intent_counts.items(), key=lambda x: x[1], reverse=True)
        return [{"intent": i[0], "count": i[1]} for i in sorted_intents[:top_n]]

    def get_satisfaction_stats(self) -> Dict:
        """Get user satisfaction statistics."""
        scores = [e.data.get("score", 0)
                 for e in self.events if e.event_type == "satisfaction_rating"]
        if not scores:
            return {"avg": 0, "total_ratings": 0, "distribution": {}}

        distribution = {}
        for s in scores:
            distribution[s] = distribution.get(s, 0) + 1

        return {
            "avg": round(sum(scores) / len(scores), 2),
            "total_ratings": len(scores),
            "distribution": distribution
        }


# =========================
# AUTOMATION SCHEDULER
# =========================

# =========================
# AUTOMATION SCHEDULER SECTION
# =========================
# The scheduler runs background jobs such as website sync, incoming PDF checks,
# timetable reminders, deadline checks, reports, and cleanup.
class AutomationScheduler:
    """Manages scheduled automation tasks."""

    def __init__(self, web_sync: WebSyncEngine,
                 timetable: TimetableAutomation,
                 notification: NotificationEngine,
                 analytics: AnalyticsEngine,
                 document_processor: Optional[DocumentProcessor] = None):
        """Initialize this component and prepare its dependencies for later use."""
        self.web_sync = web_sync
        self.timetable = timetable
        self.notification = notification
        self.analytics = analytics
        self.document_processor = document_processor
        self.tasks: Dict[str, AutomationTask] = {}
        self.running = False
        self._setup_tasks()

    def _setup_tasks(self):
        """Setup default automation tasks."""
        default_tasks = [
            AutomationTask(
                task_id="daily_timetable_reminder",
                name="Daily Timetable Reminders",
                cron_schedule="0 7 * * 1-5",  # 7 AM weekdays
            ),
            AutomationTask(
                task_id="website_sync",
                name="Website Data Synchronization",
                cron_schedule="0 2 * * *",  # 2 AM daily
            ),
            AutomationTask(
                task_id="incoming_pdf_watcher",
                name="Incoming PDF Folder Watcher",
                cron_schedule="* * * * *",  # Check incoming_pdfs every minute
            ),
            AutomationTask(
                task_id="assignment_deadline_check",
                name="Assignment Deadline Alerts",
                cron_schedule="0 8 * * *",  # 8 AM daily
            ),
            AutomationTask(
                task_id="exam_reminder_check",
                name="Exam Schedule Reminders",
                cron_schedule="0 9 * * *",  # 9 AM daily
            ),
            AutomationTask(
                task_id="analytics_report",
                name="Daily Analytics Report",
                cron_schedule="0 23 * * *",  # 11 PM daily
            ),
            AutomationTask(
                task_id="data_cleanup",
                name="Data Cleanup & Optimization",
                cron_schedule="0 3 * * 0",  # 3 AM Sundays
            ),
        ]

        for task in default_tasks:
            self.tasks[task.task_id] = task
            self._calculate_next_run(task)

    def _calculate_next_run(self, task: AutomationTask):
        """Calculate next run time based on cron schedule."""
        try:
            cron = croniter(task.cron_schedule, datetime.now())
            task.next_run = cron.get_next(datetime)
        except Exception:
            task.next_run = datetime.now() + timedelta(hours=24)

    def execute_task(self, task_id: str):
        """Execute a specific automation task."""
        if task_id not in self.tasks:
            logger.error(f"Task not found: {task_id}")
            return

        task = self.tasks[task_id]
        logger.info(f"Executing task: {task.name}")

        try:
            if task_id == "daily_timetable_reminder":
                self._run_timetable_reminders()
            elif task_id == "website_sync":
                self._run_website_sync()
            elif task_id == "incoming_pdf_watcher":
                self._run_incoming_pdf_watcher()
            elif task_id == "assignment_deadline_check":
                self._run_deadline_check()
            elif task_id == "exam_reminder_check":
                self._run_exam_reminders()
            elif task_id == "analytics_report":
                self._run_analytics_report()
            elif task_id == "data_cleanup":
                self._run_data_cleanup()

            task.last_run = datetime.now()
            task.run_count += 1
            task.status = "active"
            self._calculate_next_run(task)

        except Exception as e:
            task.error_count += 1
            task.last_error = str(e)
            task.status = "error"
            logger.error(f"Task {task_id} failed: {e}")

    def _run_timetable_reminders(self):
        """Send timetable reminders to all registered students."""
        # In production, load students from database
        students = self._get_registered_students()
        self.timetable.send_daily_reminders(students)
        self.analytics.track_event("automation_run",
                                   data={"task": "timetable_reminders",
                                         "students_notified": len(students)})

    def _run_website_sync(self):
        """Synchronize website data."""
        result = self.web_sync.incremental_sync()
        self.analytics.track_event("automation_run",
                                   data={"task": "website_sync", "result": result})

    def _run_incoming_pdf_watcher(self):
        """Append any newly dropped PDFs into the document vector store."""
        if not self.document_processor:
            return
        result = self.document_processor.sync_incoming_pdfs()
        self.analytics.track_event(
            "automation_run",
            data={"task": "incoming_pdf_watcher", "result": result}
        )

    def _run_deadline_check(self):
        """Check for upcoming assignment deadlines."""
        # In production, fetch from LMS/Aagahi integration
        self.analytics.track_event("automation_run",
                                   data={"task": "deadline_check"})

    def _run_exam_reminders(self):
        """Send exam reminders."""
        self.analytics.track_event("automation_run",
                                   data={"task": "exam_reminders"})

    def _run_analytics_report(self):
        """Generate daily analytics report."""
        stats = self.analytics.get_daily_stats(days=1)
        self.analytics.track_event("automation_run",
                                   data={"task": "analytics_report", "stats": stats})

    def _run_data_cleanup(self):
        """Clean up old data and optimize storage."""
        self.analytics.track_event("automation_run",
                                   data={"task": "data_cleanup"})

    def _get_registered_students(self) -> List[Student]:
        """Get registered students from storage."""
        if os.path.exists(Config.STUDENTS_PATH):
            try:
                with open(Config.STUDENTS_PATH, 'r') as f:
                    data = json.load(f)
                    return [Student(**s) for s in data]
            except Exception:
                pass
        return []

    def check_and_run_due_tasks(self):
        """Check for and execute any due tasks."""
        now = datetime.now()
        for task_id, task in self.tasks.items():
            if task.status == "active" and task.next_run and task.next_run <= now:
                self.execute_task(task_id)

    def start_scheduler(self):
        """Start the background scheduler."""
        self.running = True

        def _scheduler_loop():
            """Continuously checks scheduled tasks while Streamlit remains open."""
            while self.running:
                self.check_and_run_due_tasks()
                time.sleep(60)  # Check every minute

        thread = threading.Thread(target=_scheduler_loop, daemon=True)
        thread.start()
        logger.info("Automation scheduler started")

    def stop_scheduler(self):
        """Stop the background scheduler."""
        self.running = False


# =========================
# RAG ENGINE (Enhanced)
# =========================

# =========================
# RAG CHATBOT ENGINE SECTION
# =========================
# RAG means Retrieval-Augmented Generation: first retrieve trusted university
# context from vector stores, then ask the LLM to answer using that context.
class RAGEngine:
    """Enhanced RAG engine with multi-source retrieval and response generation."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.llm = ChatOpenAI(model=Config.OPENAI_MODEL, temperature=0.1, streaming=True)  # Streams GPT answers token-by-token into Streamlit.
        self.embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)
        self.bedrock = BedrockClient()

    def generate_response(self, query: str, context: str,
                         conversation_history: str = "",
                         tone: str = "Helpful & Friendly",
                         use_bedrock_fallback: bool = True) -> str:
        """Generate response using RAG with optional Bedrock fallback."""

        system_prompt = f"""You are an expert AI assistant for Salim Habib University (SHU), Karachi, Pakistan.
Tone: {tone}

CONVERSATION HISTORY:
{conversation_history}

KNOWLEDGE SOURCES:
{context}

RULES:
1. Use the provided context to answer accurately.
2. For room/lab questions, use exact room names from the timetable.
3. If information is not in context, say so clearly and provide contact info.
4. Be helpful, accurate, and concise.
5. Never invent facts not in the context.

Question: {query}
Answer:"""

        try:
            response = self.llm.invoke(system_prompt)
            return response.content
        except Exception as e:
            logger.error(f"OpenAI generation failed: {e}")

            # Fallback to Bedrock
            if use_bedrock_fallback and self.bedrock.client:
                bedrock_response = self.bedrock.generate_response(system_prompt)
                if bedrock_response:
                    return bedrock_response

            return "I'm experiencing technical difficulties. Please try again or contact SHU directly."

    def stream_response(self, query: str, context: str,
                       conversation_history: str = "",
                       tone: str = "Helpful & Friendly"):
        """Stream response tokens for real-time display."""

        system_prompt = f"""You are an expert AI assistant for Salim Habib University (SHU), Karachi, Pakistan.
Tone: {tone}

CONVERSATION HISTORY:
{conversation_history}

KNOWLEDGE SOURCES:
{context}

RULES:
1. Use the provided context to answer accurately.
2. For room/lab questions, use exact room names from the timetable.
3. If information is not in context, say so clearly.
4. Never invent facts not in the context.

Question: {query}
Answer:"""

        prompt = ChatPromptTemplate.from_template("{input}")
        chain = prompt | self.llm

        for chunk in chain.stream({"input": system_prompt}):
            if hasattr(chunk, 'content') and chunk.content:
                yield chunk.content


# =========================
# ADMIN DASHBOARD
# =========================

# =========================
# ADMIN DASHBOARD SECTION
# =========================
# The admin dashboard renders operational charts and task controls inside Streamlit.
class AdminDashboard:
    """Admin dashboard for monitoring and analytics."""

    def __init__(self, analytics: AnalyticsEngine,
                 scheduler: AutomationScheduler,
                 chat_history: ChatHistoryManager):
        """Receive dashboard data sources used to render admin charts and tables."""
        self.analytics = analytics
        self.scheduler = scheduler
        self.chat_history = chat_history

    def render(self):
        """Render the admin dashboard in Streamlit."""
        st.markdown("## 📊 Admin Dashboard")
        st.markdown("---")

        # Top-level metrics
        col1, col2, col3, col4 = st.columns(4)
        stats = self.analytics.get_daily_stats(days=7)
        total_queries = sum(d.get("queries", 0) for d in stats.values())
        total_sessions = sum(d.get("sessions", 0) for d in stats.values())
        total_users = sum(d.get("unique_users", 0) for d in stats.values())
        response_stats = self.analytics.get_response_time_stats()

        with col1:
            st.metric("Total Queries (7d)", total_queries)
        with col2:
            st.metric("Sessions (7d)", total_sessions)
        with col3:
            st.metric("Unique Users (7d)", total_users)
        with col4:
            st.metric("Avg Response Time", f"{response_stats['avg']}s")

        st.markdown("---")

        # Charts
        tab1, tab2, tab3, tab4 = st.tabs([
            "📈 Usage Trends", "🏢 Department Distribution",
            "⚙️ Automation Tasks", "💬 Recent Sessions"
        ])

        with tab1:
            self._render_usage_chart(stats)

        with tab2:
            self._render_department_chart()

        with tab3:
            self._render_automation_status()

        with tab4:
            self._render_recent_sessions()

    def _render_usage_chart(self, stats: Dict):
        """Render usage trend chart."""
        if not stats:
            st.info("No data available yet.")
            return

        df = pd.DataFrame([
            {"Date": day, "Queries": data["queries"],
             "Sessions": data["sessions"], "Users": data["unique_users"]}
            for day, data in sorted(stats.items())
        ])

        chart_theme = _get_plotly_theme()
        fig = px.line(df, x="Date", y=["Queries", "Sessions", "Users"],
                     title="Usage Trends (Last 7 Days)",
                     template=chart_theme["template"])
        fig.update_layout(
            plot_bgcolor=chart_theme["plot_bgcolor"],
            paper_bgcolor=chart_theme["paper_bgcolor"],
            font_color=chart_theme["font_color"],
        )
        st.plotly_chart(fig, use_container_width=True)

    def _render_department_chart(self):
        """Render department distribution pie chart."""
        dept_dist = self.analytics.get_department_distribution()
        if not dept_dist:
            st.info("No routing data available yet.")
            return

        chart_theme = _get_plotly_theme()
        fig = px.pie(
            values=list(dept_dist.values()),
            names=list(dept_dist.keys()),
            title="Query Distribution by Department",
            template=chart_theme["template"]
        )
        fig.update_layout(
            plot_bgcolor=chart_theme["plot_bgcolor"],
            paper_bgcolor=chart_theme["paper_bgcolor"],
            font_color=chart_theme["font_color"],
        )
        st.plotly_chart(fig, use_container_width=True)

    def _render_automation_status(self):
        """Render automation task status."""
        tasks_data = []
        for tid, task in self.scheduler.tasks.items():
            tasks_data.append({
                "Task": task.name,
                "Status": "✅ Active" if task.status == "active" else "❌ Error",
                "Last Run": task.last_run.strftime("%Y-%m-%d %H:%M") if task.last_run else "Never",
                "Next Run": task.next_run.strftime("%Y-%m-%d %H:%M") if task.next_run else "N/A",
                "Run Count": task.run_count,
                "Errors": task.error_count,
            })

        df = pd.DataFrame(tasks_data)
        st.dataframe(df, use_container_width=True, hide_index=True)

        # Manual trigger buttons
        st.markdown("### Manual Triggers")
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("🔄 Sync Website Now"):
                self.scheduler.execute_task("website_sync")
                st.success("Website sync triggered!")
        with col2:
            if st.button("📅 Send Reminders Now"):
                self.scheduler.execute_task("daily_timetable_reminder")
                st.success("Reminders sent!")
        with col3:
            if st.button("📊 Generate Report"):
                self.scheduler.execute_task("analytics_report")
                st.success("Report generated!")

    def _render_recent_sessions(self):
        """Render recent chat sessions."""
        sessions = list(self.chat_history.sessions.values())[-20:]
        sessions.reverse()

        for session in sessions:
            with st.expander(
                f"Session {session.session_id[:8]} — "
                f"{session.started_at.strftime('%Y-%m-%d %H:%M')} — "
                f"{len(session.messages)} messages"
            ):
                for msg in session.messages[-5:]:
                    role_icon = "🧑‍🎓" if msg['role'] == 'user' else "🤖"
                    st.markdown(f"{role_icon} **{msg['role'].title()}:** {msg['content'][:200]}")
                if session.department_routed:
                    st.caption(f"Routed to: {session.department_routed}")


# =========================
# MAIN APPLICATION
# =========================

# =========================
# APPLICATION ORCHESTRATION SECTION
# =========================
# This class wires together notifications, documents, web crawling, analytics,
# scheduling, chat history, and the RAG engine into one Streamlit app.
class SHUChatbotApp:
    """Main application class orchestrating all components."""

    def __init__(self):
        """Initialize this component and prepare its dependencies for later use."""
        self.notification_engine = NotificationEngine()
        self.document_processor = DocumentProcessor()
        self.web_sync = WebSyncEngine()
        self.chat_history = ChatHistoryManager()
        self.analytics = AnalyticsEngine()
        self.rag_engine = RAGEngine()
        self.query_router = QueryRouter(self.rag_engine.bedrock)
        self.timetable = TimetableAutomation(self.notification_engine)
        self.scheduler = AutomationScheduler(
            self.web_sync, self.timetable,
            self.notification_engine, self.analytics,
            self.document_processor
        )
        self.admin_dashboard = AdminDashboard(
            self.analytics, self.scheduler, self.chat_history
        )

    def initialize(self):
        """Initialize all components."""
        # Start scheduler
        if not st.session_state.get("scheduler_started"):
            self.scheduler.start_scheduler()
            st.session_state["scheduler_started"] = True

        # Initialize KB
        if not st.session_state.get("kb_ready"):
            self._init_knowledge_base()

        # Load web index
        if not st.session_state.get("web_indexed"):
            self._load_web_index()

        # Load document index so previous and newly appended PDFs are searchable.
        if not st.session_state.get("doc_indexed"):
            self._load_document_index()

    def _init_knowledge_base(self):
        """Initialize the knowledge base."""
        try:
            embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)
            os.makedirs(Config.CHROMA_KB_DIR, exist_ok=True)

            # Check if exists
            if os.path.exists(Config.CHROMA_KB_DIR):
                vs = Chroma(persist_directory=Config.CHROMA_KB_DIR,
                           embedding_function=embeddings)
                if vs._collection.count() > 0:
                    st.session_state["kb_retriever"] = vs.as_retriever(search_kwargs={"k": 10})
                    st.session_state["kb_vs"] = vs
                    st.session_state["kb_ready"] = True
                    return

            # Build new
            from langchain_core.documents import Document as LCDoc
            splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
            doc = LCDoc(page_content=SHU_MANUAL_KB,
                       metadata={"source": "SHU SP2026 Timetable + Manual KB"})
            chunks = splitter.split_documents([doc])
            vs = Chroma.from_documents(chunks, embedding=embeddings,
                                      persist_directory=Config.CHROMA_KB_DIR)
            st.session_state["kb_retriever"] = vs.as_retriever(search_kwargs={"k": 10})
            st.session_state["kb_vs"] = vs
            st.session_state["kb_ready"] = True

        except Exception as e:
            logger.error(f"KB initialization failed: {e}")

    def _load_web_index(self):
        """Load existing web index."""
        try:
            if os.path.exists(Config.CHROMA_WEB_DIR):
                embeddings = OpenAIEmbeddings(model=Config.EMBEDDING_MODEL)
                vs = Chroma(persist_directory=Config.CHROMA_WEB_DIR,
                           embedding_function=embeddings)
                if vs._collection.count() > 0:
                    st.session_state["web_retriever"] = vs.as_retriever(search_kwargs={"k": 12})
                    st.session_state["web_vs"] = vs
                    st.session_state["web_indexed"] = True
        except Exception as e:
            logger.error(f"Web index load failed: {e}")

    def _load_document_index(self):
        """Load existing uploaded document index."""
        try:
            retriever = self.document_processor.get_retriever()
            if retriever:
                st.session_state["retriever"] = retriever
                st.session_state["doc_indexed"] = True
        except Exception as e:
            logger.error(f"Document index load failed: {e}")

    def handle_query(self, query: str, session_id: str) -> str:
        """Handle a user query through the full pipeline."""
        start_time = time.time()

        # 1. Route query
        routing = self.query_router.route_query(query)
        self.analytics.track_event("query_routed", data=routing)

        # 2. Get context from all sources
        context = self._get_combined_context(query)

        # 3. Get conversation history
        conversation_history = self.chat_history.get_session_context(session_id)

        # 4. Add routing info to context
        if routing.get('contact_info'):
            context += f"\n\nDEPARTMENT CONTACT: {json.dumps(routing['contact_info'])}"

        # 5. Generate response
        tone = st.session_state.get("tone", "Helpful & Friendly")

        # 6. Track analytics
        response_time = time.time() - start_time
        self.analytics.track_event("query", data={
            "query": query[:100],
            "department": routing.get('department'),
            "intent": routing.get('intent'),
        })
        self.analytics.track_event("response_generated", data={
            "response_time": response_time
        })

        # 7. Save to chat history
        self.chat_history.add_message(session_id, "user", query)

        return context

    def _get_combined_context(self, query: str) -> str:
        """Get combined context from all knowledge sources."""
        all_docs = []

        # Direct fact lookup
        direct = get_direct_shu_fact(query)
        if direct:
            all_docs.append(Document(
                page_content=direct,
                metadata={"source": "Direct Fact", "title": "SHU Fact"}
            ))

        # KB retrieval
        if st.session_state.get("kb_retriever"):
            try:
                kb_docs = st.session_state["kb_retriever"].invoke(query)
                all_docs.extend(kb_docs)
            except Exception:
                pass

        # Web retrieval
        if st.session_state.get("web_retriever"):
            try:
                web_docs = st.session_state["web_retriever"].invoke(query)
                all_docs.extend(web_docs)
            except Exception:
                pass

        # Uploaded document retrieval. This Chroma collection is append-only:
        # manual uploads and incoming_pdfs additions share the same vector store.
        if st.session_state.get("retriever"):
            try:
                doc_docs = st.session_state["retriever"].invoke(query)
                all_docs.extend(doc_docs)
            except Exception:
                pass

        # Deduplicate
        seen, unique = set(), []
        for d in all_docs:
            key = d.page_content[:200]
            if key not in seen:
                seen.add(key)
                unique.append(d)

        # Format
        parts = []
        for d in unique:
            m = d.metadata
            loc = m.get("title", m.get("source", "Section"))
            parts.append(f"[{loc}]:\n{d.page_content}")

        return "\n\n---\n\n".join(parts)


# =========================
# MANUAL KNOWLEDGE BASE (Same as original - abbreviated for space)
# =========================

SHU_MANUAL_KB = """
=== SHU SPRING 2026 KNOWLEDGE BASE ===
Source: Official SP_2026_Semester_Wise_Timetable.pdf
Generated by aSc Timetables on 13/02/2026
University: Salim Habib University (SHU), Karachi, Pakistan

--- ROOM FLOOR / CLASS LOCATION GUIDE ---
Rooms 101, 103, 104, 105 = First Floor lecture rooms.
Rooms 201, 202, 205, 206, 207, 208, 210, 211, 212 = Second Floor rooms/labs.
Rooms 301, 309, 310 = Third Floor rooms/labs.
Room 404 = Fourth Floor lecture room.
Lec Hall 210 (Lab) = Second Floor computer lab.
Lec Hall 211 (Lab) = Second Floor computer lab.
Lec Hall 309 (Lab) = Third Floor computer lab.
Lec Hall 310 (Lab) = Third Floor computer lab.

--- TEACHER DIRECTORY ---
Prof. Dr. Sheikh Muhammad Munaf = Professor and Dean FOIT
Mr. Syed Zahid Badshah = Professor of Practice and Incharge CS Department
Mr. Ahsan Ul Haq = Lecturer/Program InCharge CTF

=== SHU GENERAL INFORMATION ===
University: Salim Habib University (SHU)
Location: Karachi, Sindh, Pakistan
Website: https://shu.edu.pk
Charter: 25 May 2015, Government of Sindh
Type: Private, Not-for-Profit
Accreditations: HEC, PEC, PCP, NCEAC, Government of Sindh
"""


def get_direct_shu_fact(query: str) -> str:
    """Exact-answer layer for common questions."""
    import re
    q = query.lower()
    q = q.replace("deann", "dean").replace("deperment", "department")

    if "dean" in q and any(w in q for w in ["cs", "computer science", "foit"]):
        return ("Prof. Dr. Sheikh Muhammad Munaf is the Professor and Dean FOIT, "
                "Faculty of Information Technology, Department of Computer Science.")

    if any(w in q for w in ["incharge", "in charge", "hod"]) and any(w in q for w in ["cs", "computer"]):
        return ("Mr. Syed Zahid Badshah is the Professor of Practice and Incharge, "
                "Department of Computer Science.")

    room_match = re.search(r"\b(\d{3})\b", q)
    if room_match and any(w in q for w in ["where", "floor", "location", "room", "class"]):
        room_no = room_match.group(1)
        floor_map = {
            "101": "First", "103": "First", "104": "First", "105": "First",
            "201": "Second", "202": "Second", "205": "Second", "206": "Second",
            "207": "Second", "208": "Second", "210": "Second", "211": "Second",
            "212": "Second", "309": "Third", "310": "Third", "404": "Fourth",
        }
        if room_no in floor_map:
            return f"Room/Lec Hall {room_no} is on the {floor_map[room_no]} Floor."

    return ""


# =========================
# STREAMLIT UI IMPLEMENTATION
# =========================

# =========================
# DARK/LIGHT MODE SECTION
# =========================
# Theme helpers keep color tokens in one place so Streamlit UI and Plotly charts
# change consistently between dark and light mode.
def _theme_tokens() -> Dict[str, str]:
    """Return the active SHU color palette from session state."""
    if st.session_state.get("theme", "dark") == "dark":
        return {
            "bg": "#050711", "bg2": "#0b1220",
            "panel": "rgba(13, 20, 35, 0.88)",
            "panel2": "rgba(18, 28, 47, 0.82)",
            "sidebar": "linear-gradient(180deg, rgba(3, 7, 18, 0.98), rgba(8, 18, 36, 0.96))",
            "text": "#eaf2ff", "muted": "#9fb0c8", "line": "rgba(119, 152, 190, 0.24)",
            "accent": "#38bdf8", "accent2": "#2563eb", "button_text": "#ffffff",
            "input": "rgba(7, 13, 26, 0.94)", "input_text": "#f8fbff",
            "table": "#0b1220", "shadow": "0 22px 70px rgba(0, 0, 0, 0.34)",
            "plotly": "plotly_dark",
        }
    return {
        "bg": "#f8fbff", "bg2": "#eaf4ff",
        "panel": "rgba(255, 255, 255, 0.92)",
        "panel2": "rgba(239, 247, 255, 0.9)",
        "sidebar": "linear-gradient(180deg, rgba(255, 255, 255, 0.98), rgba(232, 243, 255, 0.96))",
        "text": "#0f172a", "muted": "#475569", "line": "rgba(30, 64, 175, 0.18)",
        "accent": "#0ea5e9", "accent2": "#0f2f66", "button_text": "#ffffff",
        "input": "#ffffff", "input_text": "#0f172a",
        "table": "#ffffff", "shadow": "0 18px 48px rgba(15, 23, 42, 0.12)",
        "plotly": "plotly_white",
    }


def _apply_theme_css():
    """Apply the persisted Streamlit theme across native and custom UI elements."""
    t = _theme_tokens()
    st.markdown(f"""
    <style>
        :root {{
            --bg: {t["bg"]};
            --panel: {t["panel"]};
            --line: {t["line"]};
            --text: {t["text"]};
            --muted: {t["muted"]};
            --accent: {t["accent"]};
            --accent-2: {t["accent2"]};
            --shadow: {t["shadow"]};
        }}

        html, body, .stApp, [data-testid="stSidebar"], .main .block-container,
        .hero-shell, .section-card, .welcome-card, .brand-card, .stat-card,
        .notification-card, [data-testid="stChatMessage"], div[data-testid="stMetric"],
        .stButton > button, .stDownloadButton > button, .stFormSubmitButton > button,
        input, textarea, [data-baseweb="select"] > div, .stTabs [data-baseweb="tab-list"],
        [data-testid="stDataFrame"], [data-testid="stTable"] {{
            transition: background 220ms ease, background-color 220ms ease,
                        color 220ms ease, border-color 220ms ease,
                        box-shadow 220ms ease, filter 220ms ease;
        }}

        .stApp {{
            color: {t["text"]};
            background:
                radial-gradient(circle at 10% 6%, rgba(56, 189, 248, 0.18), transparent 30%),
                radial-gradient(circle at 90% 4%, rgba(37, 99, 235, 0.14), transparent 34%),
                linear-gradient(135deg, {t["bg"]} 0%, {t["bg2"]} 100%) !important;
        }}

        [data-testid="stSidebar"] {{
            background: {t["sidebar"]} !important;
            border-right: 1px solid {t["line"]};
        }}

        .hero-shell, .section-card, .welcome-card, .brand-card,
        .stat-card, .notification-card, [data-testid="stChatMessage"],
        div[data-testid="stMetric"], [data-testid="stExpander"] {{
            background: {t["panel"]} !important;
            color: {t["text"]} !important;
            border-color: {t["line"]} !important;
            box-shadow: {t["shadow"]};
        }}

        .main-title {{
            background: linear-gradient(92deg, {t["text"]} 0%, {t["accent"]} 48%, {t["accent2"]} 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }}

        h1, h2, h3, h4, h5, h6, p, label, span, div, .stMarkdown,
        [data-testid="stMarkdownContainer"], [data-testid="stCaptionContainer"] {{
            color: inherit;
        }}

        .brand-title, .welcome-title, div[data-testid="stMetricValue"],
        div[data-testid="stMetricDelta"] {{
            color: {t["text"]} !important;
        }}

        .brand-subtitle, .welcome-copy, .stat-label,
        div[data-testid="stMetric"] label, [data-testid="stCaptionContainer"] {{
            color: {t["muted"]} !important;
        }}

        .pill, .status-row, .stRadio label, .stTabs [data-baseweb="tab-list"] {{
            background: {t["panel2"]} !important;
            border-color: {t["line"]} !important;
            color: {t["text"]} !important;
        }}

        .stButton > button, .stDownloadButton > button, .stFormSubmitButton > button {{
            background: linear-gradient(135deg, {t["accent"]}, {t["accent2"]}) !important;
            color: {t["button_text"]} !important;
            border: 1px solid {t["accent"]} !important;
        }}

        .stTextInput input, .stTextArea textarea, [data-testid="stChatInput"] textarea,
        .stNumberInput input, .stDateInput input, .stTimeInput input,
        .stSelectbox div[data-baseweb="select"] > div,
        .stMultiSelect div[data-baseweb="select"] > div {{
            background: {t["input"]} !important;
            color: {t["input_text"]} !important;
            border-color: {t["line"]} !important;
        }}

        input::placeholder, textarea::placeholder {{
            color: {t["muted"]} !important;
            opacity: 0.9;
        }}

        [data-testid="stFileUploader"] section,
        [data-testid="stDataFrame"], [data-testid="stTable"],
        .stDataFrame, [data-testid="stTable"] {{
            background: {t["table"]} !important;
            color: {t["text"]} !important;
            border-color: {t["line"]} !important;
        }}

        .stTabs [data-baseweb="tab"] {{
            color: {t["muted"]} !important;
        }}

        .stTabs [aria-selected="true"] {{
            background: linear-gradient(135deg, rgba(14, 165, 233, 0.18), rgba(15, 47, 102, 0.18)) !important;
            color: {t["text"]} !important;
        }}

        [data-testid="stChatInput"] {{
            background: linear-gradient(180deg, transparent, {t["bg"]}) !important;
            border-top: 1px solid {t["line"]};
        }}

        div[data-testid="stAlert"] {{
            border-color: {t["line"]};
        }}
    </style>
    """, unsafe_allow_html=True)


def _get_plotly_theme() -> Dict[str, str]:
    """Theme Plotly charts to match the active Streamlit mode."""
    t = _theme_tokens()
    return {
        "template": t["plotly"],
        "plot_bgcolor": t["panel"],
        "paper_bgcolor": t["panel"],
        "font_color": t["text"],
    }


# =========================
# ADMIN AUTHENTICATION SECTION
# =========================
# Admin helpers protect the dashboard with credentials loaded from .env or secrets.
def _admin_credentials_configured() -> bool:
    """Return True only when both admin username and password are configured."""
    return bool(Config.ADMIN_USERNAME and Config.ADMIN_PASSWORD)


def _admin_is_authenticated() -> bool:
    """Read Streamlit session_state to check whether this browser is logged in."""
    return bool(st.session_state.get("admin_authenticated"))


def _render_admin_login() -> bool:
    """Render the admin login gate and update session_state on success."""
    st.markdown("## Admin Login")
    st.caption("Configure ADMIN_USERNAME and ADMIN_PASSWORD in .env or Streamlit secrets.")

    if not _admin_credentials_configured():
        st.error("Admin credentials are not configured. Set ADMIN_USERNAME and ADMIN_PASSWORD.")
        return False

    with st.form("admin_login_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        submitted = st.form_submit_button("Login", use_container_width=True)

    if submitted:
        valid_user = hmac.compare_digest(username, Config.ADMIN_USERNAME)
        valid_password = hmac.compare_digest(password, Config.ADMIN_PASSWORD)
        if valid_user and valid_password:
            st.session_state["admin_authenticated"] = True
            st.session_state["admin_username"] = username
            st.rerun()
        else:
            st.error("Invalid admin username or password.")

    return False


# =========================
# STREAMLIT UI SECTION
# =========================
# main() builds the complete Streamlit page: configuration, session_state,
# sidebar controls, upload tools, navigation, and page rendering.
def main():
    """Main Streamlit application entry point."""

    # Configure the browser tab and Streamlit layout before rendering widgets.
    st.set_page_config(
        page_title="SHU AI Assistant - Automation System",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    # Premium Streamlit UI theme
    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');

        :root {
            --bg: #070913;
            --panel: rgba(17, 24, 39, 0.72);
            --line: rgba(148, 163, 184, 0.22);
            --text: #e5eefc;
            --muted: #94a3b8;
            --accent: #28d7ff;
            --accent-2: #9d7cff;
            --accent-3: #4ade80;
            --warning: #fbbf24;
            --radius: 18px;
            --shadow: 0 24px 80px rgba(0, 0, 0, 0.35);
        }

        html, body, [class*="css"] {
            font-family: 'Inter', system-ui, -apple-system, BlinkMacSystemFont, sans-serif;
        }

        .stApp {
            color: var(--text);
            background:
                radial-gradient(circle at 12% 8%, rgba(40, 215, 255, 0.22), transparent 30%),
                radial-gradient(circle at 88% 4%, rgba(157, 124, 255, 0.18), transparent 34%),
                linear-gradient(135deg, #070913 0%, #0d1324 42%, #0f172a 100%);
            background-attachment: fixed;
        }

        .stApp::before {
            content: "";
            position: fixed;
            inset: 0;
            pointer-events: none;
            background-image:
                linear-gradient(rgba(255,255,255,0.035) 1px, transparent 1px),
                linear-gradient(90deg, rgba(255,255,255,0.035) 1px, transparent 1px);
            background-size: 46px 46px;
            mask-image: linear-gradient(to bottom, rgba(0,0,0,0.55), transparent 70%);
            z-index: 0;
        }

        .main .block-container {
            max-width: 1320px;
            padding: 2.2rem 2.4rem 7rem;
            position: relative;
            z-index: 1;
        }

        [data-testid="stSidebar"] {
            background:
                linear-gradient(180deg, rgba(9, 14, 27, 0.94), rgba(15, 23, 42, 0.9)),
                radial-gradient(circle at 30% 0%, rgba(40, 215, 255, 0.16), transparent 38%);
            border-right: 1px solid var(--line);
            backdrop-filter: blur(22px);
        }

        [data-testid="stSidebar"] > div:first-child {
            padding-top: 1.4rem;
        }

        [data-testid="stSidebar"] h2 {
            display: none;
        }

        [data-testid="stSidebar"] hr,
        .stMarkdown hr {
            border-color: rgba(148, 163, 184, 0.16);
            margin: 1.1rem 0;
        }

        .brand-card {
            padding: 18px 16px;
            border: 1px solid rgba(40, 215, 255, 0.22);
            border-radius: var(--radius);
            background: linear-gradient(145deg, rgba(40, 215, 255, 0.13), rgba(157, 124, 255, 0.09));
            box-shadow: 0 18px 60px rgba(0, 0, 0, 0.24);
            margin-bottom: 0.4rem;
        }

        .brand-kicker {
            color: var(--accent);
            font-size: 0.72rem;
            letter-spacing: 0.08em;
            text-transform: uppercase;
            font-weight: 800;
            margin-bottom: 6px;
        }

        .brand-title {
            color: #f8fbff;
            font-size: 1.28rem;
            font-weight: 800;
            line-height: 1.18;
        }

        .brand-subtitle {
            color: var(--muted);
            font-size: 0.82rem;
            margin-top: 8px;
        }

        .main-title {
            background: linear-gradient(92deg, #f8fbff 0%, #85e8ff 45%, #b9a7ff 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            font-size: clamp(2.1rem, 5vw, 4.5rem);
            font-weight: 800;
            line-height: 1;
            letter-spacing: 0;
            margin: 0.25rem 0 0.7rem;
        }

        .hero-shell {
            border: 1px solid rgba(148, 163, 184, 0.18);
            border-radius: 28px;
            padding: clamp(22px, 4vw, 42px);
            background:
                linear-gradient(135deg, rgba(15, 23, 42, 0.82), rgba(17, 24, 39, 0.56)),
                radial-gradient(circle at 92% 14%, rgba(40, 215, 255, 0.18), transparent 28%);
            box-shadow: var(--shadow);
            backdrop-filter: blur(24px);
            margin-bottom: 1.4rem;
            animation: riseIn 0.7s ease-out both;
        }

        .hero-meta {
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin-top: 18px;
        }

        .pill {
            display: inline-flex;
            align-items: center;
            gap: 8px;
            border-radius: 999px;
            border: 1px solid rgba(148, 163, 184, 0.22);
            color: #dbeafe;
            background: rgba(255, 255, 255, 0.06);
            padding: 8px 12px;
            font-size: 0.82rem;
            font-weight: 700;
        }

        .section-card {
            border: 1px solid rgba(148, 163, 184, 0.18);
            border-radius: var(--radius);
            padding: 18px;
            background: rgba(15, 23, 42, 0.58);
            backdrop-filter: blur(18px);
            box-shadow: 0 14px 50px rgba(0, 0, 0, 0.24);
            transition: transform 180ms ease, border-color 180ms ease, background 180ms ease;
        }

        .section-card:hover,
        .stat-card:hover,
        .notification-card:hover {
            transform: translateY(-3px);
            border-color: rgba(40, 215, 255, 0.42);
            background: rgba(15, 23, 42, 0.76);
        }

        .stat-card {
            background: linear-gradient(145deg, rgba(255,255,255,0.08), rgba(255,255,255,0.035));
            border: 1px solid rgba(148, 163, 184, 0.18);
            border-radius: 16px;
            padding: 14px 12px;
            text-align: center;
            min-height: 84px;
            transition: transform 180ms ease, border-color 180ms ease;
        }

        .stat-number {
            font-size: 1.45rem;
            font-weight: 800;
            color: var(--accent);
        }

        .stat-label {
            font-size: 0.74rem;
            color: var(--muted);
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0.05em;
        }

        .status-row {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 8px;
            padding: 9px 10px;
            margin: 6px 0;
            border-radius: 12px;
            border: 1px solid rgba(148, 163, 184, 0.12);
            background: rgba(255, 255, 255, 0.045);
            font-size: 0.88rem;
        }

        .status-dot {
            width: 9px;
            height: 9px;
            border-radius: 999px;
            display: inline-block;
            box-shadow: 0 0 18px currentColor;
        }

        .status-on { color: var(--accent-3); background: var(--accent-3); }
        .status-off { color: var(--warning); background: var(--warning); }

        .automation-badge {
            display: inline-block;
            padding: 6px 12px;
            border-radius: 999px;
            font-size: 0.75rem;
            font-weight: 800;
            background: rgba(74, 222, 128, 0.12);
            color: #86efac;
            border: 1px solid rgba(74, 222, 128, 0.28);
        }

        .notification-card {
            background: rgba(15, 23, 42, 0.68);
            border: 1px solid rgba(148, 163, 184, 0.16);
            border-left: 3px solid var(--accent);
            padding: 14px 16px;
            border-radius: 14px;
            margin: 10px 0;
            transition: transform 180ms ease, border-color 180ms ease;
        }

        .welcome-card {
            background:
                linear-gradient(135deg, rgba(40, 215, 255, 0.13), rgba(157, 124, 255, 0.11)),
                rgba(15, 23, 42, 0.68);
            border: 1px solid rgba(40, 215, 255, 0.24);
            border-radius: 22px;
            padding: 22px 24px;
            margin: 18px 0;
            box-shadow: 0 18px 70px rgba(0, 0, 0, 0.24);
            animation: riseIn 0.75s ease-out both;
        }

        .welcome-title {
            color: #ffffff;
            font-size: 1.08rem;
            font-weight: 800;
            margin-bottom: 8px;
        }

        .welcome-copy {
            color: #b9c7da;
            font-size: 0.95rem;
            line-height: 1.6;
        }

        div[data-testid="stMetric"] {
            background: linear-gradient(145deg, rgba(255,255,255,0.08), rgba(255,255,255,0.035));
            border: 1px solid rgba(148, 163, 184, 0.18);
            border-radius: 18px;
            padding: 16px;
            box-shadow: 0 16px 50px rgba(0, 0, 0, 0.2);
        }

        div[data-testid="stMetric"] label {
            color: #b6c4d8 !important;
            font-weight: 700;
        }

        div[data-testid="stMetricValue"] {
            color: #eff6ff;
            font-weight: 800;
        }

        .stButton > button,
        .stDownloadButton > button,
        .stFormSubmitButton > button {
            border: 1px solid rgba(40, 215, 255, 0.28);
            border-radius: 14px;
            color: #f8fbff;
            background: linear-gradient(135deg, rgba(40, 215, 255, 0.9), rgba(157, 124, 255, 0.85));
            box-shadow: 0 12px 34px rgba(40, 215, 255, 0.18);
            font-weight: 800;
            transition: transform 160ms ease, filter 160ms ease, box-shadow 160ms ease;
        }

        .stButton > button:hover,
        .stDownloadButton > button:hover,
        .stFormSubmitButton > button:hover {
            transform: translateY(-2px);
            filter: brightness(1.08);
            box-shadow: 0 18px 44px rgba(157, 124, 255, 0.28);
        }

        .stRadio [role="radiogroup"] {
            gap: 8px;
        }

        .stRadio label {
            border: 1px solid rgba(148, 163, 184, 0.14);
            border-radius: 14px;
            padding: 9px 12px;
            background: rgba(255,255,255,0.04);
            transition: background 160ms ease, border-color 160ms ease;
        }

        .stRadio label:hover {
            background: rgba(40, 215, 255, 0.1);
            border-color: rgba(40, 215, 255, 0.3);
        }

        .stSelectbox div[data-baseweb="select"] > div,
        .stTextInput input,
        .stTextArea textarea,
        .stMultiSelect div[data-baseweb="select"] > div {
            border-radius: 14px;
            background: rgba(255, 255, 255, 0.07);
            border-color: rgba(148, 163, 184, 0.22);
        }

        [data-testid="stFileUploader"] section {
            border-radius: 18px;
            border: 1px dashed rgba(40, 215, 255, 0.42);
            background: rgba(40, 215, 255, 0.055);
            transition: border-color 180ms ease, background 180ms ease;
        }

        [data-testid="stFileUploader"] section:hover {
            border-color: rgba(157, 124, 255, 0.65);
            background: rgba(157, 124, 255, 0.08);
        }

        .stTabs [data-baseweb="tab-list"] {
            gap: 8px;
            background: rgba(15, 23, 42, 0.48);
            border: 1px solid rgba(148, 163, 184, 0.14);
            border-radius: 16px;
            padding: 6px;
        }

        .stTabs [data-baseweb="tab"] {
            border-radius: 12px;
            color: #cbd5e1;
            font-weight: 800;
        }

        .stTabs [aria-selected="true"] {
            background: linear-gradient(135deg, rgba(40, 215, 255, 0.18), rgba(157, 124, 255, 0.18));
            color: #ffffff;
        }

        [data-testid="stChatMessage"] {
            border: 1px solid rgba(148, 163, 184, 0.16);
            border-radius: 20px;
            background: rgba(15, 23, 42, 0.62);
            backdrop-filter: blur(16px);
            box-shadow: 0 16px 50px rgba(0, 0, 0, 0.18);
            margin: 12px 0;
            animation: riseIn 0.42s ease-out both;
        }

        [data-testid="stChatInput"] {
            border-top: 1px solid rgba(148, 163, 184, 0.16);
            background: linear-gradient(180deg, transparent, rgba(7, 9, 19, 0.82));
        }

        [data-testid="stChatInput"] textarea {
            border-radius: 18px;
            border: 1px solid rgba(40, 215, 255, 0.26);
            background: rgba(15, 23, 42, 0.92);
            box-shadow: 0 16px 50px rgba(0, 0, 0, 0.28);
        }

        .stDataFrame,
        [data-testid="stTable"] {
            border-radius: 18px;
            overflow: hidden;
            border: 1px solid rgba(148, 163, 184, 0.16);
        }

        div[data-testid="stAlert"] {
            border-radius: 16px;
            border: 1px solid rgba(148, 163, 184, 0.18);
        }

        @media (prefers-color-scheme: light) {
            .stApp {
                color: #111827;
                background:
                    radial-gradient(circle at 12% 8%, rgba(40, 215, 255, 0.22), transparent 30%),
                    radial-gradient(circle at 88% 4%, rgba(157, 124, 255, 0.17), transparent 34%),
                    linear-gradient(135deg, #f8fbff 0%, #eef6ff 48%, #f7f3ff 100%);
            }

            .hero-shell,
            .section-card,
            .welcome-card,
            [data-testid="stChatMessage"],
            div[data-testid="stMetric"] {
                background: rgba(255, 255, 255, 0.72);
                color: #111827;
                box-shadow: 0 18px 60px rgba(15, 23, 42, 0.12);
            }

            .welcome-title,
            div[data-testid="stMetricValue"] {
                color: #0f172a;
            }

            .welcome-copy,
            .brand-subtitle {
                color: #475569;
            }
        }

        @media (max-width: 760px) {
            .main .block-container {
                padding: 1.2rem 1rem 6rem;
            }

            .hero-meta {
                display: grid;
                grid-template-columns: 1fr;
            }

            .pill {
                justify-content: center;
            }
        }

        @keyframes riseIn {
            from {
                opacity: 0;
                transform: translateY(12px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
    </style>
    """, unsafe_allow_html=True)

    # Initialize session state
    session_defaults = {
        "chat_history": [], "total_questions": 0,
        "session_start": datetime.now().strftime("%H:%M"),
        "current_session_id": None, "app_mode": "chat",
        "kb_ready": False, "web_indexed": False,
        "doc_indexed": False,
        "tone": "Helpful & Friendly", "scheduler_started": False,
        "pdf_processed": False, "retriever": None,
        "theme": "dark", "admin_authenticated": False,
        "admin_username": "",
    }
    for k, v in session_defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v  # session_state keeps values alive across Streamlit reruns.

    _apply_theme_css()

    # Verify API key
    if not Config.OPENAI_API_KEY:
        st.error("⚠️ Missing OPENAI_API_KEY in environment variables.")
        st.stop()

    # Initialize application
    app = SHUChatbotApp()
    app.initialize()

    # Create session if needed
    if not st.session_state["current_session_id"]:
        st.session_state["current_session_id"] = app.chat_history.create_session()

    # =========================
    # Sidebar contains controls that should remain available on every page.
    with st.sidebar:
        st.markdown("""
        <div class="brand-card">
            <div class="brand-kicker"></div>
            <div class="brand-title">SHU AI Assistant</div>
            <div class="brand-subtitle">Intelligent automation system v2.0</div>
        </div>
        """, unsafe_allow_html=True)
        light_mode = st.toggle(
            "Light mode",
            value=st.session_state.get("theme", "dark") == "light",
            key="theme_light_toggle",
        )
        selected_theme = "light" if light_mode else "dark"
        if selected_theme != st.session_state["theme"]:
            st.session_state["theme"] = selected_theme
            st.rerun()

        if _admin_is_authenticated():
            st.caption(f"Admin: {st.session_state.get('admin_username', 'authenticated')}")
            if st.button("Logout Admin", use_container_width=True):
                st.session_state["admin_authenticated"] = False
                st.session_state["admin_username"] = ""
                st.rerun()
        st.markdown("<h2 style='color:#4f8ef7;'>🎓 SHU AI Assistant</h2>",
                   unsafe_allow_html=True)
        st.caption("Intelligent Automation System v2.0")
        st.markdown("---")

        # App Mode Selection
        mode = st.radio(
            "Mode",
            ["💬 Chat", "📊 Admin Dashboard", "⚙️ Automation", "🔔 Notifications"],
            label_visibility="collapsed"
        )
        st.session_state["app_mode"] = mode

        st.markdown("---")

        # Stats
        c1, c2 = st.columns(2)
        with c1:
            st.markdown(
                f"<div class='stat-card'><div class='stat-number'>"
                f"{st.session_state['total_questions']}</div>"
                f"<div class='stat-label'>Questions</div></div>",
                unsafe_allow_html=True)
        with c2:
            st.markdown(
                f"<div class='stat-card'><div class='stat-number'>"
                f"{st.session_state['session_start']}</div>"
                f"<div class='stat-label'>Session</div></div>",
                unsafe_allow_html=True)

        st.markdown("---")

        # System Status
        st.subheader("System Status")
        statuses = [
            ("📚 Knowledge Base", st.session_state.get("kb_ready", False)),
            ("🌐 Web Index", st.session_state.get("web_indexed", False)),
            ("⚙️ Scheduler", st.session_state.get("scheduler_started", False)),
            ("🔗 Bedrock", app.rag_engine.bedrock.client is not None),
        ]
        statuses.insert(2, ("Documents", st.session_state.get("doc_indexed", False)))
        for label, active in statuses:
            icon = "✅" if active else "⚠️"
            st.markdown(f"{icon} {label}")

        st.markdown("---")

        # Tone selector
        st.session_state["tone"] = st.selectbox(
            "Response Tone",
            ["Helpful & Friendly", "Formal & Academic",
             "Simple & Brief", "Detailed & Thorough"]
        )

        st.markdown("---")

        # Document upload
        st.subheader("📄 Upload Document")
        uploaded_file = st.file_uploader(
            "Upload", type=["pdf", "docx", "doc", "pptx", "ppt"],
            label_visibility="collapsed"
        )

        if uploaded_file and st.button("Process Document", use_container_width=True):
            with st.spinner("Processing document..."):
                ext = uploaded_file.name.rsplit(".", 1)[-1].lower()
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                    tmp.write(uploaded_file.getbuffer())
                    tmp_path = tmp.name

                result = app.document_processor.process_document(
                    tmp_path, uploaded_file.name
                )
                if result["retriever"]:
                    st.session_state["retriever"] = result["retriever"]
                    st.session_state["doc_indexed"] = True
                st.session_state["pdf_processed"] = not result.get("skipped", False)
                st.session_state["doc_name"] = uploaded_file.name
                os.unlink(tmp_path)

            st.success(f"✅ Processed: {result['num_chunks']} chunks in {result['processing_time']}s")
            app.analytics.track_event("document_processed", data={
                "file": uploaded_file.name, "chunks": result['num_chunks']
            })

        if st.button("Sync New PDFs", use_container_width=True):
            with st.spinner(f"Checking {Config.INCOMING_PDFS_DIR}..."):
                result = app.document_processor.sync_incoming_pdfs()
                retriever = app.document_processor.get_retriever()
                if retriever:
                    st.session_state["retriever"] = retriever
                    st.session_state["doc_indexed"] = True

            processed_count = len(result["processed"])
            skipped_count = len(result["skipped"])
            failed_count = len(result["failed"])
            if failed_count:
                st.error(f"Processed {processed_count}, skipped {skipped_count}, failed {failed_count}.")
            elif processed_count:
                chunks = sum(item["chunks"] for item in result["processed"])
                st.success(f"Synced {processed_count} PDF(s), added {chunks} chunks.")
            else:
                st.info(f"No new PDFs found. Skipped {skipped_count} already processed file(s).")
            app.analytics.track_event("incoming_pdf_sync", data=result)

        # Web crawl button
        st.markdown("---")
        if not st.session_state.get("web_indexed"):
            if st.button("🌐 Crawl shu.edu.pk", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()

                def progress_cb(current, total, docs_count):
                    """Update the Streamlit progress bar while website crawling runs."""
                    progress_bar.progress(min(int(current / total * 100), 99))
                    status_text.text(f"Crawling... {current}/{total} | {docs_count} pages")

                result = app.web_sync.full_sync(progress_callback=progress_cb)
                if result["status"] == "success":
                    st.session_state["web_retriever"] = result["retriever"]
                    st.session_state["web_vs"] = result["vector_store"]
                    st.session_state["web_indexed"] = True
                    progress_bar.progress(100)
                    st.success(f"✅ {result['pages_scraped']} pages indexed!")
                    st.rerun()

        # Clear chat
        if st.button("🗑️ Clear Chat", use_container_width=True):
            st.session_state["chat_history"] = []
            st.session_state["total_questions"] = 0
            st.rerun()

    # =========================

    if mode == "📊 Admin Dashboard":
        if _admin_is_authenticated():
            app.admin_dashboard.render()
        else:
            _render_admin_login()

    elif mode == "⚙️ Automation":
        _render_automation_page(app)

    elif mode == "🔔 Notifications":
        _render_notifications_page(app)

    else:  # Chat mode
        _render_chat_page(app)


def _render_chat_page(app: SHUChatbotApp):
    """Render the main chat interface."""
    st.markdown("""
    <div class="hero-shell">
        <div class="pill">Live university knowledge engine</div>
        <div class="main-title">Salim Habib University AI Assistant</div>
        <div class="welcome-copy">
            Premium AI support for timetable intelligence, document search,
            university operations, reminders, analytics, and automation.
        </div>
        <div class="hero-meta">
            <span class="pill">SP 2026 timetable</span>
            <span class="pill">Rooms and labs</span>
            <span class="pill">Teacher routing</span>
            <span class="pill">Automation ready</span>
        </div>
    </div>
    """, unsafe_allow_html=True)
    st.caption("SP 2026 Timetable · Rooms · Labs · Teachers · Programs · Automation ⚡")

    # Welcome message
    if not st.session_state["chat_history"]:
        st.markdown("""
        <div class="welcome-card">
            <div class="welcome-title">Ask anything across SHU knowledge, documents, and schedules.</div>
            <div class="welcome-copy">
                The assistant has the Spring 2026 official timetable loaded with rooms,
                labs, teachers, and schedules. Upload documents in the sidebar or start
                with one of the quick prompts below.
            </div>
        </div>
        """, unsafe_allow_html=True)

        # Suggestions
        suggestions = [
            "Where is the OS Lab?",
            "Who teaches AI?",
            "CS Semester 5 subjects?",
            "Dean of CS department?",
            "Show my schedule",
            "Available scholarships?"
        ]
        cols = st.columns(3)
        for i, s in enumerate(suggestions):
            with cols[i % 3]:
                if st.button(s, key=f"sug_{i}", use_container_width=True):
                    st.session_state["suggested_click"] = s
                    st.rerun()

    # Display chat history
    for role, msg in st.session_state["chat_history"]:
        if role == "user":
            with st.chat_message("user", avatar="🧑‍🎓"):
                st.markdown(msg)
        elif role == "bot":
            with st.chat_message("assistant", avatar="🤖"):
                st.markdown(msg)
        elif role == "routing":
            with st.expander("🔀 Query Routing Info", expanded=False):
                st.markdown(msg)

    # Chat input
    query = st.chat_input("Ask about rooms, labs, teachers, subjects, schedules...")  # Chat input captures the student question.

    if st.session_state.get("suggested_click"):
        query = st.session_state["suggested_click"]
        st.session_state["suggested_click"] = None

    if query:
        st.session_state["total_questions"] += 1
        st.session_state["chat_history"].append(("user", query))

        with st.chat_message("user", avatar="🧑‍🎓"):
            st.markdown(query)

        # Route query
        routing = app.query_router.route_query(query)

        # Get context
        context = app.handle_query(query, st.session_state["current_session_id"])

        # Generate streaming response
        with st.chat_message("assistant", avatar="🤖"):
            tone = st.session_state.get("tone", "Helpful & Friendly")
            conversation_history = app.chat_history.get_session_context(
                st.session_state["current_session_id"]
            )

            answer = st.write_stream(
                app.rag_engine.stream_response(
                    query, context, conversation_history, tone
                )
            )

            # Show routing badge
            dept = routing.get('department', 'General')
            st.caption(f"GPT-4o · Routed: {dept} · {datetime.now().strftime('%H:%M:%S')}")

        # Save response
        app.chat_history.add_message(
            st.session_state["current_session_id"], "assistant", answer
        )
        st.session_state["chat_history"].append(("bot", answer))

        # Show routing info
        if routing.get('confidence', 0) > 0.6:
            routing_info = (
                f"**Department:** {routing['department']}\n"
                f"**Intent:** {routing.get('intent', 'general')}\n"
                f"**Priority:** {routing.get('priority', 'MEDIUM')}\n"
                f"**Confidence:** {routing.get('confidence', 0):.0%}"
            )
            if routing.get('contact_info'):
                routing_info += f"\n**Contact:** {routing['contact_info'].get('email', '')}"
            st.session_state["chat_history"].append(("routing", routing_info))

        # Track analytics
        app.analytics.track_event("query", data={
            "query": query[:100],
            "department": routing.get('department'),
            "intent": routing.get('intent')
        })

        st.rerun()


def _render_automation_page(app: SHUChatbotApp):
    """Render automation management page."""
    st.markdown("## ⚙️ Automation Management")
    st.markdown("---")

    # Task status table
    st.subheader("Scheduled Tasks")
    tasks_data = []
    for tid, task in app.scheduler.tasks.items():
        tasks_data.append({
            "ID": tid,
            "Task": task.name,
            "Schedule": task.cron_schedule,
            "Status": task.status,
            "Last Run": task.last_run.strftime("%m/%d %H:%M") if task.last_run else "Never",
            "Next Run": task.next_run.strftime("%m/%d %H:%M") if task.next_run else "N/A",
            "Runs": task.run_count,
            "Errors": task.error_count,
        })

    df = pd.DataFrame(tasks_data)
    st.dataframe(df, use_container_width=True, hide_index=True)

    st.markdown("---")

    # Manual triggers
    st.subheader("Manual Task Execution")
    col1, col2, col3 = st.columns(3)

    with col1:
        st.markdown("#### 🌐 Website Sync")
        st.write("Crawl and index shu.edu.pk")
        if st.button("Execute Website Sync", use_container_width=True):
            with st.spinner("Syncing..."):
                app.scheduler.execute_task("website_sync")
            st.success("✅ Website sync completed!")

    with col2:
        st.markdown("#### 📅 Timetable Reminders")
        st.write("Send today's schedule to students")
        if st.button("Send Reminders", use_container_width=True):
            app.scheduler.execute_task("daily_timetable_reminder")
            st.success("✅ Reminders sent!")

    with col3:
        st.markdown("#### 📊 Analytics Report")
        st.write("Generate usage report")
        if st.button("Generate Report", use_container_width=True):
            app.scheduler.execute_task("analytics_report")
            st.success("✅ Report generated!")

    st.markdown("---")

    # Add custom automation
    st.subheader("Add Custom Automation")
    with st.form("add_automation"):
        task_name = st.text_input("Task Name")
        cron_expr = st.text_input("Cron Expression", placeholder="0 8 * * 1-5")
        task_type = st.selectbox("Type", [
            "Send Notification", "Sync Data", "Generate Report", "Custom Script"
        ])
        submit = st.form_submit_button("Add Task")

        if submit and task_name and cron_expr:
            st.success(f"✅ Task '{task_name}' added with schedule: {cron_expr}")


def _render_notifications_page(app: SHUChatbotApp):
    """Render notifications management page."""
    st.markdown("## 🔔 Notification Center")
    st.markdown("---")

    tab1, tab2, tab3 = st.tabs(["📤 Send Notification", "📋 History", "⚙️ Settings"])

    with tab1:
        st.subheader("Send University Announcement")
        with st.form("send_notification"):
            title = st.text_input("Title")
            message = st.text_area("Message")
            notif_type = st.selectbox("Type", [
                "Announcement", "Timetable Change", "Exam Update",
                "Assignment Reminder", "System Update"
            ])
            target = st.multiselect("Target Groups", [
                "All Students", "CS Students", "SE Students",
                "AI Students", "Semester 1", "Semester 2",
                "Semester 3", "Semester 4", "Semester 5",
                "Semester 6", "Semester 7", "Semester 8"
            ])
            send_email = st.checkbox("Send via Email", value=True)
            submit = st.form_submit_button("Send Notification")

            if submit and title and message:
                app.notification_engine.create_notification(
                    type=NotificationType.ANNOUNCEMENT,
                    title=title,
                    message=message,
                    recipient_id="all" if "All Students" in target else ",".join(target)
                )
                st.success(f"✅ Notification sent: {title}")
                app.analytics.track_event("notification_sent", data={
                    "title": title, "type": notif_type, "targets": target
                })

    with tab2:
        st.subheader("Recent Notifications")
        recent = app.notification_engine.notifications[-20:]
        recent.reverse()

        for notif in recent:
            status_icon = "✅" if notif.sent else "⏳"
            st.markdown(f"""
            <div class='notification-card'>
                {status_icon} <b>{notif.title}</b><br>
                <small>{notif.type.value} | {notif.created_at.strftime('%Y-%m-%d %H:%M')}</small><br>
                {notif.message[:100]}...
            </div>
            """, unsafe_allow_html=True)

    with tab3:
        st.subheader("Notification Settings")
        st.markdown("#### Default Notification Schedule")

        col1, col2 = st.columns(2)
        with col1:
            st.checkbox("Daily Timetable Reminders", value=True)
            st.checkbox("Assignment Deadline Alerts (3 days before)", value=True)
            st.checkbox("Exam Reminders (7 days before)", value=True)
        with col2:
            st.checkbox("University Announcements", value=True)
            st.checkbox("System Updates", value=False)
            st.checkbox("Weekly Summary Email", value=True)

        st.markdown("#### Email Configuration")
        st.text_input("SMTP Host", value=Config.SMTP_HOST, disabled=True)
        st.text_input("From Email", value=Config.EMAIL_FROM, disabled=True)
        st.caption("Configure in .env file for security")


# =========================
# ENTRY POINT
# =========================

if __name__ == "__main__":
    main()
